"""
Universal Document Creator - FastAPI Backend
AI-powered document generation with multi-model support
Real authentication, SQLite database, PDF engine, brand style extraction
"""

import os
from pathlib import Path
from dotenv import load_dotenv
load_dotenv(Path(__file__).resolve().parent.parent / ".env")

from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Depends, Request, Body
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from fastapi.responses import Response, FileResponse, HTMLResponse
from sse_starlette.sse import EventSourceResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field
from typing import List, Optional, Dict, Any
import json
import re
import base64
import asyncio
import aiofiles
import sqlite3
import hashlib
import secrets
import jwt
import stripe
from datetime import datetime, timedelta
import io
from io import BytesIO
from concurrent.futures import ThreadPoolExecutor

# Thread pool for blocking calls (Ollama, PDF, etc.)
_executor = ThreadPoolExecutor(max_workers=4)
_stream_executor = ThreadPoolExecutor(max_workers=4)

# PDF generation
from xhtml2pdf import pisa

# Gemini (new SDK)
from google import genai
from google.genai import types as genai_types

# Ollama (connects to Docker container on port 11434)
from ollama import Client as OllamaClientClass
ollama_client = OllamaClientClass(host="http://localhost:11434", timeout=120)

# Image processing
from PIL import Image

app = FastAPI(
    title="Universal Document Creator API",
    description="AI-powered document generation with multi-model support",
    version="2.0.0"
)

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configuration
SKILLS_FILE = os.path.join(os.path.dirname(__file__), "skills.md")
UPLOADED_SKILLS_DIR = os.path.join(os.path.dirname(__file__), "uploaded_skills")
UPLOADS_DIR = os.path.join(os.path.dirname(__file__), "uploads")
DB_PATH = os.path.join(os.path.dirname(__file__), "udc.db")
OUTPUT_DIR = "C:/UniversalDoc_Output"
os.makedirs(OUTPUT_DIR, exist_ok=True)
JWT_SECRET = os.getenv("JWT_SECRET", secrets.token_hex(32))
JWT_ALGORITHM = "HS256"
JWT_EXPIRY_HOURS = 72

# Database URL (empty = SQLite, postgresql://... = PostgreSQL)
DATABASE_URL = os.getenv("DATABASE_URL", "")

# API Keys (loaded from environment or set via settings endpoint)
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY", "")

# Stripe
STRIPE_SECRET_KEY = os.getenv("STRIPE_SECRET_KEY", "")
STRIPE_WEBHOOK_SECRET = os.getenv("STRIPE_WEBHOOK_SECRET", "")

# Stripe price mapping (plan_name -> monthly price in cents)
STRIPE_PLAN_PRICES = {
    "pro": 1900,        # $19/mo
    "enterprise": 4900,  # $49/mo
}

# Ensure directories exist
os.makedirs(UPLOADED_SKILLS_DIR, exist_ok=True)
os.makedirs(UPLOADS_DIR, exist_ok=True)

# Security
security = HTTPBearer(auto_error=False)

# ==================== MODEL DEFINITIONS ====================

AVAILABLE_MODELS = [
    {
        "id": "gemini-3-flash-preview",
        "name": "Gemini 3 Flash",
        "provider": "google",
        "description": "Google's frontier model — fast, pro-grade reasoning, vision capable",
        "vision": True,
        "censored": True,
        "tag": "Recommended"
    },
    {
        "id": "gemini-2.0-flash",
        "name": "Gemini 2.0 Flash",
        "provider": "google",
        "description": "Google's previous gen — fast, high quality, vision capable",
        "vision": True,
        "censored": True,
        "tag": ""
    },
    {
        "id": "qwen2.5vl:7b",
        "name": "Qwen Vision (Unrestricted)",
        "provider": "ollama",
        "description": "Local AI — vision capable, no content restrictions",
        "vision": True,
        "censored": False,
        "tag": "Unrestricted"
    },
    {
        "id": "dolphin3:8b",
        "name": "Dolphin 3 (Fast)",
        "provider": "ollama",
        "description": "Local AI — fast text generation, no content restrictions",
        "vision": False,
        "censored": False,
        "tag": "Fast"
    },
    {
        "id": "claude-sonnet-4-20250514",
        "name": "Claude Sonnet 4",
        "provider": "anthropic",
        "description": "Anthropic's best all-round model — excellent writing, reasoning, analysis",
        "vision": True,
        "censored": True,
        "tag": "Premium"
    },
    {
        "id": "claude-haiku-4-5-20251001",
        "name": "Claude Haiku 4.5",
        "provider": "anthropic",
        "description": "Anthropic's fastest model — great for quick drafts and simple docs",
        "vision": True,
        "censored": True,
        "tag": "Fast"
    },
]

# Image generation models (separate from text models)
IMAGE_MODELS = [
    {
        "id": "gemini-3-pro-image-preview",
        "name": "Gemini 3 Pro Image",
        "provider": "google",
        "description": "Best quality — reasons about design then generates stunning visuals",
        "tag": "Recommended"
    },
    {
        "id": "gemini-3.1-flash-image-preview",
        "name": "Gemini 3.1 Flash Image",
        "provider": "google",
        "description": "Latest Flash — fast with great quality, smart design reasoning",
        "tag": "Latest"
    },
    {
        "id": "gemini-2.5-flash-image",
        "name": "Gemini 2.5 Flash Image",
        "provider": "google",
        "description": "Stable image generation with text+image mixed output",
        "tag": ""
    },
    {
        "id": "imagen-4.0-generate-001",
        "name": "Imagen 4",
        "provider": "google",
        "description": "Dedicated image generator — photorealistic, no text reasoning",
        "tag": ""
    },
    {
        "id": "imagen-4.0-fast-generate-001",
        "name": "Imagen 4 Fast",
        "provider": "google",
        "description": "Fastest image generation — good for quick drafts",
        "tag": "Fast"
    },
]

# ==================== DATABASE ====================
# SQLite by default for desktop use. For PostgreSQL (production deploy),
# run backend/migrate_to_postgres.py to copy data, then set DATABASE_URL
# in .env. Full PostgreSQL query layer is a planned future enhancement.

def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("PRAGMA foreign_keys=ON")
    return conn

def init_db():
    conn = get_db()
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            email TEXT UNIQUE NOT NULL,
            name TEXT NOT NULL,
            password_hash TEXT NOT NULL,
            salt TEXT NOT NULL,
            plan TEXT NOT NULL DEFAULT 'free',
            created_at TEXT NOT NULL DEFAULT (datetime('now')),
            updated_at TEXT NOT NULL DEFAULT (datetime('now'))
        );

        CREATE TABLE IF NOT EXISTS documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            title TEXT NOT NULL DEFAULT 'Untitled Document',
            content TEXT NOT NULL,
            html_content TEXT,
            skill_used TEXT,
            prompt TEXT,
            model_used TEXT,
            brand_style TEXT,
            created_at TEXT NOT NULL DEFAULT (datetime('now')),
            updated_at TEXT NOT NULL DEFAULT (datetime('now')),
            FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE CASCADE
        );

        CREATE TABLE IF NOT EXISTS brand_profiles (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            name TEXT NOT NULL DEFAULT 'My Brand',
            primary_color TEXT DEFAULT '#2563eb',
            secondary_color TEXT DEFAULT '#64748b',
            accent_color TEXT DEFAULT '#f59e0b',
            font_family TEXT DEFAULT 'Arial, sans-serif',
            logo_url TEXT,
            screenshot_path TEXT,
            style_json TEXT,
            created_at TEXT NOT NULL DEFAULT (datetime('now')),
            FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE CASCADE
        );

        CREATE TABLE IF NOT EXISTS settings (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            key TEXT UNIQUE NOT NULL,
            value TEXT NOT NULL
        );

        CREATE TABLE IF NOT EXISTS teams (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL,
            owner_id INTEGER NOT NULL,
            plan TEXT NOT NULL DEFAULT 'enterprise',
            max_members INTEGER DEFAULT 50,
            created_at TEXT NOT NULL DEFAULT (datetime('now')),
            FOREIGN KEY (owner_id) REFERENCES users(id) ON DELETE CASCADE
        );

        CREATE TABLE IF NOT EXISTS team_members (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            team_id INTEGER NOT NULL,
            user_id INTEGER NOT NULL,
            role TEXT NOT NULL DEFAULT 'member',
            joined_at TEXT NOT NULL DEFAULT (datetime('now')),
            FOREIGN KEY (team_id) REFERENCES teams(id) ON DELETE CASCADE,
            FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE CASCADE,
            UNIQUE(team_id, user_id)
        );

        CREATE TABLE IF NOT EXISTS shared_templates (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            team_id INTEGER NOT NULL,
            name TEXT NOT NULL,
            content TEXT NOT NULL,
            skill_name TEXT,
            created_by INTEGER NOT NULL,
            created_at TEXT NOT NULL DEFAULT (datetime('now')),
            FOREIGN KEY (team_id) REFERENCES teams(id) ON DELETE CASCADE,
            FOREIGN KEY (created_by) REFERENCES users(id) ON DELETE CASCADE
        );

        CREATE TABLE IF NOT EXISTS audit_logs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            action TEXT NOT NULL,
            details TEXT,
            ip_address TEXT,
            created_at TEXT NOT NULL DEFAULT (datetime('now')),
            FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE SET NULL
        );

        CREATE INDEX IF NOT EXISTS idx_users_email ON users(email);
        CREATE INDEX IF NOT EXISTS idx_documents_user_id ON documents(user_id);
        CREATE INDEX IF NOT EXISTS idx_brand_profiles_user_id ON brand_profiles(user_id);
        CREATE INDEX IF NOT EXISTS idx_team_members_team ON team_members(team_id);
        CREATE INDEX IF NOT EXISTS idx_team_members_user ON team_members(user_id);
        CREATE INDEX IF NOT EXISTS idx_shared_templates_team ON shared_templates(team_id);
        CREATE INDEX IF NOT EXISTS idx_audit_logs_user ON audit_logs(user_id);
        CREATE INDEX IF NOT EXISTS idx_audit_logs_action ON audit_logs(action);
        CREATE INDEX IF NOT EXISTS idx_audit_logs_created ON audit_logs(created_at);

        CREATE TABLE IF NOT EXISTS document_versions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER NOT NULL,
            content TEXT NOT NULL,
            version_number INTEGER NOT NULL DEFAULT 1,
            created_by INTEGER,
            created_at TEXT NOT NULL DEFAULT (datetime('now')),
            FOREIGN KEY (document_id) REFERENCES documents(id) ON DELETE CASCADE,
            FOREIGN KEY (created_by) REFERENCES users(id) ON DELETE SET NULL
        );
        CREATE INDEX IF NOT EXISTS idx_doc_versions_doc ON document_versions(document_id);

        CREATE TABLE IF NOT EXISTS signature_requests (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER NOT NULL,
            sender_id INTEGER NOT NULL,
            recipient_email TEXT NOT NULL,
            recipient_name TEXT,
            status TEXT NOT NULL DEFAULT 'pending',
            sign_token TEXT UNIQUE NOT NULL,
            signed_at TEXT,
            signature_data TEXT,
            ip_address TEXT,
            created_at TEXT NOT NULL DEFAULT (datetime('now')),
            FOREIGN KEY (document_id) REFERENCES documents(id) ON DELETE CASCADE,
            FOREIGN KEY (sender_id) REFERENCES users(id) ON DELETE CASCADE
        );
        CREATE INDEX IF NOT EXISTS idx_sig_requests_doc ON signature_requests(document_id);
        CREATE INDEX IF NOT EXISTS idx_sig_requests_token ON signature_requests(sign_token);

        CREATE TABLE IF NOT EXISTS document_views (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER NOT NULL,
            viewer_token TEXT,
            viewer_ip TEXT,
            viewer_email TEXT,
            duration_seconds INTEGER DEFAULT 0,
            scroll_depth REAL DEFAULT 0,
            opened_at TEXT NOT NULL DEFAULT (datetime('now')),
            last_activity TEXT,
            FOREIGN KEY (document_id) REFERENCES documents(id) ON DELETE CASCADE
        );
        CREATE INDEX IF NOT EXISTS idx_doc_views_doc ON document_views(document_id);

        CREATE TABLE IF NOT EXISTS voice_profiles (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            name TEXT NOT NULL DEFAULT 'My Voice',
            tone TEXT,
            formality TEXT,
            vocabulary_level TEXT,
            avg_sentence_length TEXT,
            common_phrases TEXT,
            writing_rules TEXT,
            raw_analysis TEXT,
            sample_count INTEGER DEFAULT 0,
            created_at TEXT NOT NULL DEFAULT (datetime('now')),
            updated_at TEXT NOT NULL DEFAULT (datetime('now')),
            FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE CASCADE
        );
        CREATE INDEX IF NOT EXISTS idx_voice_profiles_user ON voice_profiles(user_id);

        CREATE TABLE IF NOT EXISTS marketplace_templates (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            author_id INTEGER NOT NULL,
            name TEXT NOT NULL,
            description TEXT,
            category TEXT DEFAULT 'general',
            skill_content TEXT NOT NULL,
            tags TEXT,
            downloads INTEGER DEFAULT 0,
            rating_sum INTEGER DEFAULT 0,
            rating_count INTEGER DEFAULT 0,
            is_approved INTEGER DEFAULT 1,
            created_at TEXT NOT NULL DEFAULT (datetime('now')),
            FOREIGN KEY (author_id) REFERENCES users(id) ON DELETE CASCADE
        );
        CREATE TABLE IF NOT EXISTS template_reviews (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            template_id INTEGER NOT NULL,
            user_id INTEGER NOT NULL,
            rating INTEGER NOT NULL CHECK(rating >= 1 AND rating <= 5),
            review TEXT,
            created_at TEXT NOT NULL DEFAULT (datetime('now')),
            FOREIGN KEY (template_id) REFERENCES marketplace_templates(id) ON DELETE CASCADE,
            FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE CASCADE,
            UNIQUE(template_id, user_id)
        );
        CREATE INDEX IF NOT EXISTS idx_marketplace_category ON marketplace_templates(category);
        CREATE INDEX IF NOT EXISTS idx_marketplace_author ON marketplace_templates(author_id);
    """)
    conn.commit()

    # Migration: add stripe_customer_id column if missing
    cursor = conn.execute("PRAGMA table_info(users)")
    columns = [row[1] for row in cursor.fetchall()]
    if "stripe_customer_id" not in columns:
        conn.execute("ALTER TABLE users ADD COLUMN stripe_customer_id TEXT")
        conn.commit()

    # Add email verification and password reset columns
    for col_sql in [
        "ALTER TABLE users ADD COLUMN email_verified INTEGER DEFAULT 1",
        "ALTER TABLE users ADD COLUMN verify_token TEXT",
        "ALTER TABLE users ADD COLUMN reset_token TEXT",
        "ALTER TABLE users ADD COLUMN reset_token_expires TEXT",
    ]:
        try:
            conn.execute(col_sql)
        except:
            pass
    conn.commit()

    # Add admin and rate limiting columns
    try:
        conn.execute("ALTER TABLE users ADD COLUMN is_admin INTEGER DEFAULT 0")
    except:
        pass
    try:
        conn.execute("ALTER TABLE users ADD COLUMN generation_count INTEGER DEFAULT 0")
    except:
        pass
    try:
        conn.execute("ALTER TABLE users ADD COLUMN generation_reset TEXT")
    except:
        pass
    conn.commit()

    # Add web publishing columns to documents
    for col_sql in ["ALTER TABLE documents ADD COLUMN is_published INTEGER DEFAULT 0",
                     "ALTER TABLE documents ADD COLUMN publish_slug TEXT",
                     "ALTER TABLE documents ADD COLUMN publish_password TEXT"]:
        try: conn.execute(col_sql)
        except: pass
    conn.commit()

    conn.close()

@app.on_event("startup")
async def startup():
    init_db()

# ==================== AUDIT LOGGING ====================

def log_audit(user_id: int | None, action: str, details: str = None, ip: str = None):
    """Log an audit event"""
    try:
        conn = get_db()
        conn.execute(
            "INSERT INTO audit_logs (user_id, action, details, ip_address) VALUES (?, ?, ?, ?)",
            (user_id, action, details, ip)
        )
        conn.commit()
        conn.close()
    except Exception:
        pass  # Don't let audit logging break the app

# ==================== RATE LIMITING ====================

def check_generation_limit(user: dict, conn) -> bool:
    """Check if user can generate. Returns True if allowed, raises HTTPException if not."""
    if not user:
        return True  # Guest users not rate limited for now

    # Admins and paid plans bypass limits
    if user.get("is_admin") or user.get("plan") in ("pro", "enterprise"):
        return True

    user_id = user["id"]
    row = conn.execute(
        "SELECT generation_count, generation_reset FROM users WHERE id = ?",
        (user_id,)
    ).fetchone()

    if not row:
        return True

    # Reset counter if month changed
    now = datetime.utcnow()
    reset_date = row["generation_reset"]
    if reset_date:
        try:
            reset_dt = datetime.fromisoformat(reset_date)
            if now.month != reset_dt.month or now.year != reset_dt.year:
                conn.execute(
                    "UPDATE users SET generation_count = 0, generation_reset = ? WHERE id = ?",
                    (now.isoformat(), user_id)
                )
                conn.commit()
                return True
        except:
            pass
    else:
        conn.execute(
            "UPDATE users SET generation_reset = ? WHERE id = ?",
            (now.isoformat(), user_id)
        )
        conn.commit()

    count = row["generation_count"] or 0
    if count >= 10:  # Free plan limit: 10 generations/month
        raise HTTPException(
            status_code=429,
            detail="Free plan limit reached (10 generations/month). Upgrade to Pro for unlimited generations."
        )

    return True

# ==================== AUTH HELPERS ====================

def hash_password(password: str, salt: str) -> str:
    return hashlib.pbkdf2_hmac(
        'sha256',
        password.encode('utf-8'),
        salt.encode('utf-8'),
        100_000
    ).hex()

def create_token(user_id: int, email: str) -> str:
    payload = {
        "sub": str(user_id),
        "email": email,
        "exp": datetime.utcnow() + timedelta(hours=JWT_EXPIRY_HOURS),
        "iat": datetime.utcnow()
    }
    return jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGORITHM)

def verify_token(token: str) -> dict:
    try:
        return jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token has expired")
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail="Invalid token")

async def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security)):
    if not credentials:
        raise HTTPException(status_code=401, detail="Not authenticated")
    payload = verify_token(credentials.credentials)
    user_id = payload.get("sub")
    conn = get_db()
    user = conn.execute("SELECT id, email, name, plan, created_at, is_admin FROM users WHERE id = ?", (user_id,)).fetchone()
    conn.close()
    if not user:
        raise HTTPException(status_code=401, detail="User not found")
    return dict(user)

async def get_optional_user(credentials: HTTPAuthorizationCredentials = Depends(security)):
    if not credentials:
        return None
    try:
        payload = verify_token(credentials.credentials)
        user_id = payload.get("sub")
        conn = get_db()
        user = conn.execute("SELECT id, email, name, plan, created_at, is_admin FROM users WHERE id = ?", (user_id,)).fetchone()
        conn.close()
        if user:
            return dict(user)
    except Exception:
        pass
    return None

# ==================== DATA MODELS ====================

class RegisterRequest(BaseModel):
    name: str
    email: str
    password: str

class LoginRequest(BaseModel):
    email: str
    password: str

class UserResponse(BaseModel):
    id: int
    email: str
    name: str
    plan: str
    created_at: str

class AuthResponse(BaseModel):
    user: UserResponse
    token: str

class SkillInput(BaseModel):
    name: str
    type: str = "string"
    description: str
    required: bool = True

class SkillDefinition(BaseModel):
    name: str
    description: str
    inputs: Dict[str, Any]
    outputs: str
    template: str

class DocumentRequest(BaseModel):
    skill_name: Optional[str] = None
    prompt: str
    parameters: Dict[str, Any] = Field(default_factory=dict)
    model: str = "gemini-3-flash-preview"
    temperature: float = 0.7
    max_tokens: int = 4000
    title: Optional[str] = None
    brand_profile_id: Optional[int] = None
    output_format: str = "markdown"  # markdown, html, pdf
    language: Optional[str] = None  # e.g., "Spanish", "French", "German"

class DocumentResponse(BaseModel):
    content: str
    html_content: Optional[str] = None
    skill_used: Optional[str] = None
    model_used: str
    export_formats: List[str] = ["text", "markdown", "html", "pdf"]
    generated_at: str
    document_id: Optional[int] = None

class SkillChainRequest(BaseModel):
    chain: List[Dict[str, Any]]
    initial_prompt: str
    model: str = "gemini-3-flash-preview"

class RefinementRequest(BaseModel):
    previous_content: str
    feedback: str
    skill_name: Optional[str] = None
    model: str = "gemini-3-flash-preview"

class DocumentUpdate(BaseModel):
    title: Optional[str] = None
    content: Optional[str] = None

class BrandProfileRequest(BaseModel):
    name: str = "My Brand"

class ApiKeyUpdate(BaseModel):
    gemini_api_key: Optional[str] = None
    anthropic_api_key: Optional[str] = None

class StripeCheckoutRequest(BaseModel):
    plan_name: str  # "pro" or "enterprise"

class StripeKeyUpdate(BaseModel):
    stripe_secret_key: Optional[str] = None
    stripe_webhook_secret: Optional[str] = None

# ==================== AI GENERATION ENGINE ====================

async def generate_with_gemini(prompt: str, model_name: str = "gemini-3-flash-preview",
                                temperature: float = 0.7,
                                max_tokens: int = 4000, image_data: bytes = None) -> str:
    """Generate using Google Gemini API (new SDK)"""
    api_key = GEMINI_API_KEY

    # Check DB for stored key if env var not set
    if not api_key:
        conn = get_db()
        row = conn.execute("SELECT value FROM settings WHERE key = 'gemini_api_key'").fetchone()
        conn.close()
        if row:
            api_key = row["value"]

    if not api_key:
        raise HTTPException(status_code=400, detail="Gemini API key not configured. Go to Settings to add it.")

    try:
        client = genai.Client(api_key=api_key)

        config = genai_types.GenerateContentConfig(
            temperature=temperature,
            max_output_tokens=max_tokens,
        )

        contents = []
        if image_data:
            # Vision request — include image
            image = Image.open(BytesIO(image_data))
            contents = [prompt, image]
        else:
            contents = prompt

        response = client.models.generate_content(
            model=model_name,
            contents=contents,
            config=config,
        )

        return response.text
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gemini error: {str(e)}")


async def generate_with_ollama(prompt: str, model_name: str = "dolphin3:8b",
                                temperature: float = 0.7, image_data: bytes = None) -> str:
    """Generate using local Ollama model (Docker container)"""
    try:
        messages = [{"role": "user", "content": prompt}]

        if image_data:
            img_b64 = base64.b64encode(image_data).decode("utf-8")
            messages = [{
                "role": "user",
                "content": prompt,
                "images": [img_b64]
            }]

        def _call_ollama():
            return ollama_client.chat(
                model=model_name,
                messages=messages,
                options={"temperature": temperature}
            )

        loop = asyncio.get_event_loop()
        response = await loop.run_in_executor(_executor, _call_ollama)
        return response.message.content
    except Exception as e:
        error_msg = str(e)
        if "not found" in error_msg.lower() or "pull" in error_msg.lower():
            raise HTTPException(
                status_code=503,
                detail=f"Model '{model_name}' not installed. Run: docker exec ollama ollama pull {model_name}"
            )
        if "connection" in error_msg.lower() or "refused" in error_msg.lower():
            raise HTTPException(
                status_code=503,
                detail="Ollama is not running. Start the Docker container: docker start ollama"
            )
        raise HTTPException(status_code=500, detail=f"Ollama error: {error_msg}")


async def generate_with_anthropic(prompt: str, model_name: str = "claude-sonnet-4-20250514",
                                   temperature: float = 0.7, max_tokens: int = 4000,
                                   image_data: bytes = None) -> str:
    """Generate using Anthropic Claude API"""
    api_key = ANTHROPIC_API_KEY
    if not api_key:
        conn = get_db()
        row = conn.execute("SELECT value FROM settings WHERE key = 'anthropic_api_key'").fetchone()
        conn.close()
        if row:
            api_key = row["value"]
    if not api_key:
        raise HTTPException(status_code=400, detail="Anthropic API key not configured. Go to Settings to add it.")

    try:
        import httpx
        headers = {
            "x-api-key": api_key,
            "anthropic-version": "2023-06-01",
            "content-type": "application/json"
        }
        messages = [{"role": "user", "content": prompt}]
        if image_data:
            img_b64 = base64.b64encode(image_data).decode("utf-8")
            messages = [{"role": "user", "content": [
                {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": img_b64}},
                {"type": "text", "text": prompt}
            ]}]

        payload = {
            "model": model_name,
            "max_tokens": max_tokens,
            "temperature": temperature,
            "messages": messages
        }

        async with httpx.AsyncClient(timeout=120) as client:
            resp = await client.post("https://api.anthropic.com/v1/messages", headers=headers, json=payload)

        if resp.status_code != 200:
            raise HTTPException(status_code=resp.status_code, detail=f"Anthropic error: {resp.text}")

        data = resp.json()
        return data["content"][0]["text"]
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Anthropic error: {str(e)}")


async def generate_image(prompt: str, model_id: str = "gemini-3-pro-image-preview",
                          aspect_ratio: str = "1:1") -> bytes:
    """Generate an image using Google Imagen or Gemini image models"""
    api_key = GEMINI_API_KEY
    if not api_key:
        raise HTTPException(status_code=400, detail="Gemini API key required for image generation")

    try:
        if model_id.startswith("imagen"):
            # Imagen API — uses :predict endpoint
            import httpx
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_id}:predict?key={api_key}"
            payload = {
                "instances": [{"prompt": prompt}],
                "parameters": {
                    "sampleCount": 1,
                    "aspectRatio": aspect_ratio,
                    "personGeneration": "allow_adult"
                }
            }
            async with httpx.AsyncClient(timeout=120) as client:
                resp = await client.post(url, json=payload)
            if resp.status_code != 200:
                raise HTTPException(status_code=resp.status_code, detail=f"Imagen error: {resp.text}")
            data = resp.json()
            return base64.b64decode(data["predictions"][0]["bytesBase64Encoded"])
        else:
            # Gemini native image gen — uses generateContent with responseModalities
            import httpx
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_id}:generateContent"
            headers = {"x-goog-api-key": api_key, "Content-Type": "application/json"}
            payload = {
                "contents": [{"parts": [{"text": prompt}]}],
                "generationConfig": {
                    "responseModalities": ["TEXT", "IMAGE"],
                    "imageConfig": {"aspectRatio": aspect_ratio}
                }
            }
            async with httpx.AsyncClient(timeout=120) as client:
                resp = await client.post(url, headers=headers, json=payload)
            if resp.status_code != 200:
                raise HTTPException(status_code=resp.status_code, detail=f"Gemini image error: {resp.text}")
            data = resp.json()
            for part in data["candidates"][0]["content"]["parts"]:
                if "inlineData" in part:
                    return base64.b64decode(part["inlineData"]["data"])
            raise HTTPException(status_code=500, detail="No image in response")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Image generation error: {str(e)}")


async def generate_content(prompt: str, model_id: str = "gemini-3-flash-preview",
                           temperature: float = 0.7, max_tokens: int = 4000,
                           image_data: bytes = None) -> str:
    """Route to the right AI provider based on model selection"""

    # Find model config
    model_config = next((m for m in AVAILABLE_MODELS if m["id"] == model_id), None)

    if not model_config:
        raise HTTPException(status_code=400, detail=f"Unknown model: {model_id}")

    provider = model_config["provider"]

    if provider == "google":
        return await generate_with_gemini(prompt, model_id, temperature, max_tokens, image_data)
    elif provider == "ollama":
        return await generate_with_ollama(prompt, model_id, temperature, image_data)
    elif provider == "anthropic":
        return await generate_with_anthropic(prompt, model_id, temperature, max_tokens, image_data)
    else:
        raise HTTPException(status_code=400, detail=f"Unknown provider: {provider}")

# ==================== PDF ENGINE ====================

DOCUMENT_CSS = """
@page {
    size: A4;
    margin: 2cm;
}
body {
    font-family: 'Helvetica', 'Arial', sans-serif;
    font-size: 11pt;
    line-height: 1.6;
    color: #1a1a1a;
    max-width: 100%;
}
h1 {
    font-size: 24pt;
    color: {primary_color};
    border-bottom: 3px solid {primary_color};
    padding-bottom: 8px;
    margin-top: 20px;
    margin-bottom: 15px;
}
h2 {
    font-size: 18pt;
    color: {secondary_color};
    margin-top: 25px;
    margin-bottom: 10px;
}
h3 {
    font-size: 14pt;
    color: #444;
    margin-top: 18px;
    margin-bottom: 8px;
}
p {
    margin: 8px 0;
    text-align: justify;
}
ul, ol {
    margin: 10px 0;
    padding-left: 25px;
}
li {
    margin: 4px 0;
}
table {
    width: 100%;
    border-collapse: collapse;
    margin: 15px 0;
}
th {
    background-color: {primary_color};
    color: white;
    padding: 10px 12px;
    text-align: left;
    font-weight: bold;
}
td {
    padding: 8px 12px;
    border-bottom: 1px solid #e0e0e0;
}
tr:nth-child(even) td {
    background-color: #f8f9fa;
}
blockquote {
    border-left: 4px solid {accent_color};
    padding: 10px 20px;
    margin: 15px 0;
    background-color: #fafafa;
    font-style: italic;
}
code {
    background-color: #f4f4f4;
    padding: 2px 6px;
    border-radius: 3px;
    font-family: 'Courier New', monospace;
    font-size: 10pt;
}
pre {
    background-color: #f4f4f4;
    padding: 15px;
    border-radius: 5px;
    overflow-x: auto;
    font-family: 'Courier New', monospace;
    font-size: 9pt;
    line-height: 1.4;
}
.header-bar {
    background-color: {primary_color};
    color: white;
    padding: 20px 25px;
    margin: -2cm -2cm 25px -2cm;
    text-align: center;
}
.header-bar h1 {
    color: white;
    border: none;
    margin: 0;
    padding: 0;
    font-size: 28pt;
}
.header-bar .subtitle {
    color: rgba(255,255,255,0.85);
    font-size: 12pt;
    margin-top: 5px;
}
.footer {
    text-align: center;
    font-size: 8pt;
    color: #999;
    margin-top: 40px;
    padding-top: 10px;
    border-top: 1px solid #e0e0e0;
}
.badge {
    display: inline-block;
    background-color: {accent_color};
    color: white;
    padding: 3px 10px;
    border-radius: 12px;
    font-size: 9pt;
    font-weight: bold;
}
strong {
    color: #111;
}
hr {
    border: none;
    border-top: 2px solid #e0e0e0;
    margin: 20px 0;
}
"""

def markdown_to_html(md_content: str) -> str:
    """Convert markdown to HTML (basic conversion)"""
    html = md_content

    # Headers (process longest first so #### doesn't match as # + ###)
    html = re.sub(r'^###### (.+)$', r'<h6>\1</h6>', html, flags=re.MULTILINE)
    html = re.sub(r'^##### (.+)$', r'<h5>\1</h5>', html, flags=re.MULTILINE)
    html = re.sub(r'^#### (.+)$', r'<h4>\1</h4>', html, flags=re.MULTILINE)
    html = re.sub(r'^### (.+)$', r'<h3>\1</h3>', html, flags=re.MULTILINE)
    html = re.sub(r'^## (.+)$', r'<h2>\1</h2>', html, flags=re.MULTILINE)
    html = re.sub(r'^# (.+)$', r'<h1>\1</h1>', html, flags=re.MULTILINE)

    # Bold and italic
    html = re.sub(r'\*\*\*(.+?)\*\*\*', r'<strong><em>\1</em></strong>', html)
    html = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', html)
    html = re.sub(r'\*(.+?)\*', r'<em>\1</em>', html)

    # Code blocks — extract content between ``` fences, escape HTML inside
    def _code_block_replace(m):
        code = m.group(1)
        code = code.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        return f'<pre><code>{code}</code></pre>'
    html = re.sub(r'```[\w]*\n(.*?)```', _code_block_replace, html, flags=re.DOTALL)
    html = re.sub(r'`([^`]+)`', r'<code>\1</code>', html)

    # Blockquotes
    html = re.sub(r'^> (.+)$', r'<blockquote>\1</blockquote>', html, flags=re.MULTILINE)

    # Horizontal rules
    html = re.sub(r'^---+$', r'<hr>', html, flags=re.MULTILINE)

    # Tables — convert markdown pipe tables to HTML before list processing
    table_lines = html.split('\n')
    table_result = []
    in_table = False
    header_done = False
    for i, line in enumerate(table_lines):
        stripped = line.strip()
        if re.match(r'^\|(.+)\|$', stripped):
            cells = [c.strip() for c in stripped.strip('|').split('|')]
            # Check if next line is a separator (|---|---|)
            is_separator = all(re.match(r'^[-:]+$', c) for c in cells)
            if is_separator:
                header_done = True
                continue
            if not in_table:
                table_result.append('<table>')
                in_table = True
                header_done = False
            if not header_done and in_table:
                # This is the header row
                table_result.append('<thead><tr>')
                for c in cells:
                    table_result.append(f'<th>{c}</th>')
                table_result.append('</tr></thead><tbody>')
            else:
                table_result.append('<tr>')
                for c in cells:
                    table_result.append(f'<td>{c}</td>')
                table_result.append('</tr>')
        else:
            if in_table:
                table_result.append('</tbody></table>')
                in_table = False
                header_done = False
            table_result.append(line)
    if in_table:
        table_result.append('</tbody></table>')
    html = '\n'.join(table_result)

    # Unordered lists
    lines = html.split('\n')
    result = []
    in_list = False
    in_ol = False
    for line in lines:
        stripped = line.strip()
        if re.match(r'^[-*] ', stripped):
            if not in_list:
                result.append('<ul>')
                in_list = True
            item = re.sub(r'^[-*] ', '', stripped)
            result.append(f'<li>{item}</li>')
        else:
            if in_list:
                result.append('</ul>')
                in_list = False
            # Ordered lists
            ol_match = re.match(r'^(\d+)\. (.+)', stripped)
            if ol_match:
                if not in_ol:
                    result.append('<ol>')
                    in_ol = True
                result.append(f'<li>{ol_match.group(2)}</li>')
            elif stripped:
                if in_ol:
                    result.append('</ol>')
                    in_ol = False
                # Wrap plain text in <p> if not already an HTML tag
                if not stripped.startswith('<'):
                    result.append(f'<p>{stripped}</p>')
                else:
                    result.append(stripped)
            else:
                if in_ol:
                    result.append('</ol>')
                    in_ol = False
                result.append('')
    if in_list:
        result.append('</ul>')
    if in_ol:
        result.append('</ol>')

    return '\n'.join(result)


def generate_pdf_bytes(html_content: str, brand_style: dict = None, include_footer: bool = True) -> bytes:
    """Generate PDF from HTML content using xhtml2pdf"""
    style = brand_style or {}
    primary = style.get("primary_color", "#2563eb")
    secondary = style.get("secondary_color", "#64748b")
    accent = style.get("accent_color", "#f59e0b")
    font = style.get("font_family", "Helvetica, Arial, sans-serif")

    css = DOCUMENT_CSS.replace("{primary_color}", primary)\
                       .replace("{secondary_color}", secondary)\
                       .replace("{accent_color}", accent)\
                       .replace("{font_family}", font)

    footer_html = '<div class="footer">Generated by Universal Document Creator</div>' if include_footer else ''

    full_html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <style>{css}</style>
</head>
<body>
{html_content}
{footer_html}
</body>
</html>"""

    buffer = BytesIO()
    pisa_status = pisa.CreatePDF(full_html, dest=buffer)

    if pisa_status.err:
        raise HTTPException(status_code=500, detail="PDF generation failed")

    return buffer.getvalue()

# ==================== BRAND STYLE EXTRACTION ====================

async def extract_brand_style(image_data: bytes, model_id: str = "gemini-3-flash-preview") -> dict:
    """Use vision AI to extract brand colors/style from a website screenshot"""
    prompt = """Analyze this website screenshot and extract the brand style. Return ONLY valid JSON with these fields:
{
    "primary_color": "#hex color of the main/primary brand color",
    "secondary_color": "#hex color of the secondary color",
    "accent_color": "#hex color used for buttons/CTAs/highlights",
    "background_color": "#hex background color",
    "text_color": "#hex main text color",
    "font_family": "detected or closest matching web-safe font family",
    "style_description": "brief description of the overall design style",
    "logo_position": "left/center/right - where the logo appears"
}
Return ONLY the JSON, no markdown formatting, no explanation."""

    result = await generate_content(prompt, model_id, temperature=0.1, image_data=image_data)

    # Parse JSON from response
    try:
        # Try to extract JSON from response
        json_match = re.search(r'\{[^{}]*\}', result, re.DOTALL)
        if json_match:
            return json.loads(json_match.group())
        return json.loads(result)
    except json.JSONDecodeError:
        # Return defaults if parsing fails
        return {
            "primary_color": "#2563eb",
            "secondary_color": "#64748b",
            "accent_color": "#f59e0b",
            "background_color": "#ffffff",
            "text_color": "#1a1a1a",
            "font_family": "Arial, sans-serif",
            "style_description": "Could not parse style from screenshot",
            "logo_position": "left"
        }

# ==================== SKILLS SYSTEM ====================

def parse_skills_from_markdown(content: str) -> List[SkillDefinition]:
    skills = []
    skill_blocks = re.split(r'### Skill:', content)

    for block in skill_blocks[1:]:
        lines = block.strip().split('\n')
        skill_data = {
            "name": "", "description": "", "inputs": {},
            "outputs": "", "template": ""
        }
        current_section = None
        template_lines = []
        in_template = False

        for line in lines:
            line = line.strip()
            if line.startswith('- **Name**:'):
                skill_data["name"] = line.split(':', 1)[1].strip()
            elif line.startswith('- **Description**:'):
                skill_data["description"] = line.split(':', 1)[1].strip()
            elif line.startswith('- **Inputs**:'):
                current_section = "inputs"
            elif line.startswith('- **Outputs**:'):
                skill_data["outputs"] = line.split(':', 1)[1].strip()
                current_section = None
            elif line.startswith('- **Template**:'):
                in_template = True
                template_part = line.split(':', 1)[1].strip()
                if template_part and template_part != '|':
                    template_lines.append(template_part)
            elif in_template:
                if line.startswith('- '):
                    in_template = False
                else:
                    template_lines.append(line.lstrip('| '))
            elif current_section == "inputs" and line.startswith('- '):
                input_match = re.match(r'-\s+(\w+)\s+\(([^)]+)\):\s*(.+)', line)
                if input_match:
                    skill_data["inputs"][input_match.group(1)] = {
                        "type": input_match.group(2),
                        "description": input_match.group(3)
                    }

        skill_data["template"] = '\n'.join(template_lines).strip()
        if skill_data["name"]:
            skills.append(SkillDefinition(**skill_data))
    return skills

def load_all_skills() -> List[SkillDefinition]:
    all_skills = []
    try:
        with open(SKILLS_FILE, 'r', encoding='utf-8') as f:
            content = f.read()
        all_skills.extend(parse_skills_from_markdown(content))
    except Exception as e:
        print(f"Error loading built-in skills: {e}")
    try:
        for filename in os.listdir(UPLOADED_SKILLS_DIR):
            if filename.endswith('.md'):
                filepath = os.path.join(UPLOADED_SKILLS_DIR, filename)
                with open(filepath, 'r', encoding='utf-8') as f:
                    content = f.read()
                all_skills.extend(parse_skills_from_markdown(content))
    except Exception as e:
        print(f"Error loading uploaded skills: {e}")
    return all_skills

def get_skill_by_name(name: str) -> Optional[SkillDefinition]:
    for skill in load_all_skills():
        if skill.name.lower() == name.lower():
            return skill
    return None

def apply_skill_template(skill: SkillDefinition, parameters: Dict[str, Any]) -> str:
    template = skill.template
    for key, value in parameters.items():
        placeholder = "{" + key + "}"
        if isinstance(value, list):
            value = ", ".join(str(v) for v in value)
        template = template.replace(placeholder, str(value))
    # Remove ICU select blocks for undefined optional params: {var, select, undefined {} other {text}}
    template = re.sub(r'\{[^}]+,\s*select,\s*undefined\s*\{\}\s*other\s*\{[^}]*\}\}', '', template)
    return template

# ==================== AUTH ENDPOINTS ====================

@app.post("/api/auth/register", response_model=AuthResponse)
async def register_user(req: RegisterRequest):
    if not req.email or not req.password or not req.name:
        raise HTTPException(status_code=400, detail="Name, email, and password are required")
    if len(req.password) < 6:
        raise HTTPException(status_code=400, detail="Password must be at least 6 characters")

    conn = get_db()
    existing = conn.execute("SELECT id FROM users WHERE email = ?", (req.email.lower(),)).fetchone()
    if existing:
        conn.close()
        raise HTTPException(status_code=409, detail="Email already registered")

    salt = secrets.token_hex(16)
    password_hash = hash_password(req.password, salt)
    cursor = conn.execute(
        "INSERT INTO users (email, name, password_hash, salt, plan) VALUES (?, ?, ?, ?, ?)",
        (req.email.lower(), req.name, password_hash, salt, "free")
    )
    conn.commit()
    user_id = cursor.lastrowid
    user = conn.execute("SELECT id, email, name, plan, created_at FROM users WHERE id = ?", (user_id,)).fetchone()
    conn.close()

    log_audit(user_id, "register", f"New account: {req.email}")
    return AuthResponse(user=UserResponse(**dict(user)), token=create_token(user_id, req.email.lower()))

@app.post("/api/auth/login", response_model=AuthResponse)
async def login_user(req: LoginRequest, request: Request):
    conn = get_db()
    user = conn.execute("SELECT * FROM users WHERE email = ?", (req.email.lower(),)).fetchone()
    conn.close()
    if not user:
        raise HTTPException(status_code=401, detail="Invalid email or password")

    if hash_password(req.password, user["salt"]) != user["password_hash"]:
        raise HTTPException(status_code=401, detail="Invalid email or password")

    log_audit(user["id"], "login", f"Login from {request.client.host if hasattr(request, 'client') else 'unknown'}")
    return AuthResponse(
        user=UserResponse(id=user["id"], email=user["email"], name=user["name"],
                          plan=user["plan"], created_at=user["created_at"]),
        token=create_token(user["id"], user["email"])
    )

@app.get("/api/auth/me", response_model=UserResponse)
async def get_me(user=Depends(get_current_user)):
    return UserResponse(**user)

# ==================== PASSWORD RESET ENDPOINTS ====================

@app.post("/api/auth/forgot-password")
async def forgot_password(request_data: dict = Body(...)):
    """Generate a password reset token"""
    email = request_data.get("email", "").strip()
    if not email:
        raise HTTPException(status_code=400, detail="Email is required")

    conn = get_db()
    user = conn.execute("SELECT id, email FROM users WHERE email = ?", (email.lower(),)).fetchone()

    if not user:
        conn.close()
        # Don't reveal if email exists - always return success
        return {"message": "If an account with that email exists, a reset link has been sent."}

    # Generate reset token
    import secrets as sec
    reset_token = sec.token_urlsafe(32)
    expires = (datetime.utcnow() + timedelta(hours=1)).isoformat()

    conn.execute(
        "UPDATE users SET reset_token = ?, reset_token_expires = ? WHERE id = ?",
        (reset_token, expires, user["id"])
    )
    conn.commit()
    conn.close()

    # In production, send email with reset link containing the token
    # For local/desktop app, we return the token directly (dev mode)
    return {
        "message": "If an account with that email exists, a reset link has been sent.",
        "reset_token": reset_token  # Remove in production - only for local dev
    }


@app.post("/api/auth/reset-password")
async def reset_password(request_data: dict = Body(...)):
    """Reset password using a valid reset token"""
    token = request_data.get("token", "").strip()
    new_password = request_data.get("new_password", "").strip()

    if not token or not new_password:
        raise HTTPException(status_code=400, detail="Token and new password are required")

    if len(new_password) < 6:
        raise HTTPException(status_code=400, detail="Password must be at least 6 characters")

    conn = get_db()
    user = conn.execute(
        "SELECT id, reset_token_expires FROM users WHERE reset_token = ?",
        (token,)
    ).fetchone()

    if not user:
        conn.close()
        raise HTTPException(status_code=400, detail="Invalid or expired reset token")

    # Check expiration
    if user["reset_token_expires"]:
        expires = datetime.fromisoformat(user["reset_token_expires"])
        if datetime.utcnow() > expires:
            conn.close()
            raise HTTPException(status_code=400, detail="Reset token has expired")

    # Update password and clear token
    import secrets as sec
    salt = sec.token_hex(16)
    password_hash_val = hash_password(new_password, salt)

    conn.execute(
        "UPDATE users SET password_hash = ?, salt = ?, reset_token = NULL, reset_token_expires = NULL WHERE id = ?",
        (password_hash_val, salt, user["id"])
    )
    conn.commit()
    conn.close()

    return {"message": "Password reset successfully. You can now sign in with your new password."}

# ==================== SETTINGS ENDPOINTS ====================

@app.post("/api/settings/api-keys")
async def update_api_keys(req: ApiKeyUpdate, user=Depends(get_current_user)):
    """Save API keys to database"""
    conn = get_db()
    if req.gemini_api_key is not None:
        conn.execute(
            "INSERT INTO settings (key, value) VALUES ('gemini_api_key', ?) "
            "ON CONFLICT(key) DO UPDATE SET value = ?",
            (req.gemini_api_key, req.gemini_api_key)
        )
    if req.anthropic_api_key is not None:
        conn.execute(
            "INSERT INTO settings (key, value) VALUES ('anthropic_api_key', ?) "
            "ON CONFLICT(key) DO UPDATE SET value = ?",
            (req.anthropic_api_key, req.anthropic_api_key)
        )
        global ANTHROPIC_API_KEY
        ANTHROPIC_API_KEY = req.anthropic_api_key
    conn.commit()
    conn.close()
    return {"message": "API keys updated"}

@app.get("/api/settings/api-keys")
async def get_api_keys(user=Depends(get_current_user)):
    """Check which API keys are configured"""
    conn = get_db()
    gemini = conn.execute("SELECT value FROM settings WHERE key = 'gemini_api_key'").fetchone()
    anthropic = conn.execute("SELECT value FROM settings WHERE key = 'anthropic_api_key'").fetchone()
    conn.close()
    return {
        "gemini_configured": bool(gemini) or bool(GEMINI_API_KEY),
        "anthropic_configured": bool(anthropic) or bool(ANTHROPIC_API_KEY),
    }

# ==================== STRIPE SETTINGS ====================

@app.post("/api/settings/stripe-key")
async def update_stripe_keys(req: StripeKeyUpdate, user=Depends(get_current_user)):
    """Save Stripe keys to database (admin-style, any authenticated user)"""
    conn = get_db()
    if req.stripe_secret_key is not None:
        conn.execute(
            "INSERT INTO settings (key, value) VALUES ('stripe_secret_key', ?) "
            "ON CONFLICT(key) DO UPDATE SET value = ?",
            (req.stripe_secret_key, req.stripe_secret_key)
        )
        global STRIPE_SECRET_KEY
        STRIPE_SECRET_KEY = req.stripe_secret_key
    if req.stripe_webhook_secret is not None:
        conn.execute(
            "INSERT INTO settings (key, value) VALUES ('stripe_webhook_secret', ?) "
            "ON CONFLICT(key) DO UPDATE SET value = ?",
            (req.stripe_webhook_secret, req.stripe_webhook_secret)
        )
        global STRIPE_WEBHOOK_SECRET
        STRIPE_WEBHOOK_SECRET = req.stripe_webhook_secret
    conn.commit()
    conn.close()
    return {"message": "Stripe keys updated"}

@app.get("/api/settings/stripe-key")
async def get_stripe_key_status(user=Depends(get_current_user)):
    """Check if Stripe keys are configured"""
    conn = get_db()
    sk = conn.execute("SELECT value FROM settings WHERE key = 'stripe_secret_key'").fetchone()
    wh = conn.execute("SELECT value FROM settings WHERE key = 'stripe_webhook_secret'").fetchone()
    conn.close()
    return {
        "stripe_configured": bool(sk) or bool(STRIPE_SECRET_KEY),
        "webhook_configured": bool(wh) or bool(STRIPE_WEBHOOK_SECRET),
    }

# ==================== STRIPE PAYMENT ENDPOINTS ====================

def _get_stripe_key() -> str:
    """Get Stripe secret key from env or DB"""
    if STRIPE_SECRET_KEY:
        return STRIPE_SECRET_KEY
    conn = get_db()
    row = conn.execute("SELECT value FROM settings WHERE key = 'stripe_secret_key'").fetchone()
    conn.close()
    return row["value"] if row else ""

def _get_webhook_secret() -> str:
    """Get Stripe webhook secret from env or DB"""
    if STRIPE_WEBHOOK_SECRET:
        return STRIPE_WEBHOOK_SECRET
    conn = get_db()
    row = conn.execute("SELECT value FROM settings WHERE key = 'stripe_webhook_secret'").fetchone()
    conn.close()
    return row["value"] if row else ""

@app.post("/api/stripe/create-checkout-session")
async def create_checkout_session(req: StripeCheckoutRequest, user=Depends(get_current_user)):
    """Create a Stripe Checkout session for Pro or Enterprise plan"""
    sk = _get_stripe_key()
    if not sk:
        raise HTTPException(status_code=400, detail="Stripe not configured. Add your Stripe secret key in Settings.")

    plan = req.plan_name.lower()
    if plan not in STRIPE_PLAN_PRICES:
        raise HTTPException(status_code=400, detail=f"Invalid plan: {plan}. Choose 'pro' or 'enterprise'.")

    stripe.api_key = sk

    try:
        # Check if user already has a Stripe customer ID
        conn = get_db()
        db_user = conn.execute("SELECT stripe_customer_id FROM users WHERE id = ?", (user["id"],)).fetchone()
        customer_id = db_user["stripe_customer_id"] if db_user and db_user["stripe_customer_id"] else None
        conn.close()

        # Create or reuse Stripe customer
        if not customer_id:
            customer = stripe.Customer.create(
                email=user["email"],
                name=user.get("name", ""),
                metadata={"user_id": str(user["id"])}
            )
            customer_id = customer.id
            conn = get_db()
            conn.execute("UPDATE users SET stripe_customer_id = ? WHERE id = ?", (customer_id, user["id"]))
            conn.commit()
            conn.close()

        # Create checkout session
        session = stripe.checkout.Session.create(
            customer=customer_id,
            payment_method_types=["card"],
            mode="subscription",
            line_items=[{
                "price_data": {
                    "currency": "usd",
                    "recurring": {"interval": "month"},
                    "product_data": {"name": f"Universal Document Creator — {plan.title()} Plan"},
                    "unit_amount": STRIPE_PLAN_PRICES[plan],
                },
                "quantity": 1,
            }],
            metadata={"user_id": str(user["id"]), "plan_name": plan},
            success_url=os.getenv("STRIPE_SUCCESS_URL", "http://localhost:5173/dashboard?payment=success"),
            cancel_url=os.getenv("STRIPE_CANCEL_URL", "http://localhost:5173/pricing?payment=cancelled"),
        )

        return {"url": session.url, "session_id": session.id}

    except stripe.error.StripeError as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/api/stripe/webhook")
async def stripe_webhook(request: Request):
    """Handle Stripe webhook events (checkout.session.completed -> update user plan)"""
    wh_secret = _get_webhook_secret()
    payload = await request.body()
    sig_header = request.headers.get("stripe-signature", "")

    sk = _get_stripe_key()
    if not sk:
        raise HTTPException(status_code=400, detail="Stripe not configured")
    stripe.api_key = sk

    # Verify webhook signature if secret is set
    if wh_secret:
        try:
            event = stripe.Webhook.construct_event(payload, sig_header, wh_secret)
        except ValueError:
            raise HTTPException(status_code=400, detail="Invalid payload")
        except stripe.error.SignatureVerificationError:
            raise HTTPException(status_code=400, detail="Invalid signature")
    else:
        # No webhook secret — parse raw (dev/testing only)
        try:
            event = json.loads(payload)
        except json.JSONDecodeError:
            raise HTTPException(status_code=400, detail="Invalid JSON payload")

    # Handle checkout completed
    if event.get("type") == "checkout.session.completed":
        session_data = event["data"]["object"]
        user_id = session_data.get("metadata", {}).get("user_id")
        plan_name = session_data.get("metadata", {}).get("plan_name")
        customer_id = session_data.get("customer")

        if user_id and plan_name:
            conn = get_db()
            conn.execute(
                "UPDATE users SET plan = ?, stripe_customer_id = ?, updated_at = datetime('now') WHERE id = ?",
                (plan_name, customer_id, int(user_id))
            )
            conn.commit()
            conn.close()

    # Handle subscription deleted (downgrade to free)
    elif event.get("type") == "customer.subscription.deleted":
        sub = event["data"]["object"]
        customer_id = sub.get("customer")
        if customer_id:
            conn = get_db()
            conn.execute(
                "UPDATE users SET plan = 'free', updated_at = datetime('now') WHERE stripe_customer_id = ?",
                (customer_id,)
            )
            conn.commit()
            conn.close()

    return {"status": "ok"}

@app.post("/api/stripe/customer-portal")
async def create_customer_portal(user=Depends(get_current_user)):
    """Create a Stripe Billing Portal session for subscription management"""
    sk = _get_stripe_key()
    if not sk:
        raise HTTPException(status_code=400, detail="Stripe not configured. Add your Stripe secret key in Settings.")

    stripe.api_key = sk

    # Get user's Stripe customer ID
    conn = get_db()
    db_user = conn.execute("SELECT stripe_customer_id FROM users WHERE id = ?", (user["id"],)).fetchone()
    conn.close()

    if not db_user or not db_user["stripe_customer_id"]:
        raise HTTPException(status_code=400, detail="No billing account found. Subscribe to a plan first.")

    try:
        portal_session = stripe.billing_portal.Session.create(
            customer=db_user["stripe_customer_id"],
            return_url=os.getenv("STRIPE_PORTAL_RETURN_URL", "http://localhost:5173/dashboard"),
        )
        return {"url": portal_session.url}
    except stripe.error.StripeError as e:
        raise HTTPException(status_code=400, detail=str(e))

# ==================== BRAND PROFILE ENDPOINTS ====================

@app.post("/api/brand/upload-screenshot")
async def upload_brand_screenshot(
    file: UploadFile = File(...),
    name: str = Form("My Brand"),
    model: str = Form("gemini-3-flash-preview"),
    user=Depends(get_current_user)
):
    """Upload a website screenshot and extract brand style"""
    if not file.content_type or not file.content_type.startswith("image/"):
        raise HTTPException(status_code=400, detail="Only image files are allowed")

    image_data = await file.read()

    # Save screenshot
    timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    ext = file.filename.split('.')[-1] if file.filename else "png"
    filename = f"brand_{user['id']}_{timestamp}.{ext}"
    filepath = os.path.join(UPLOADS_DIR, filename)
    async with aiofiles.open(filepath, 'wb') as f:
        await f.write(image_data)

    # Extract brand style using vision AI
    style = await extract_brand_style(image_data, model)

    # Save to database
    conn = get_db()
    cursor = conn.execute(
        "INSERT INTO brand_profiles (user_id, name, primary_color, secondary_color, accent_color, "
        "font_family, screenshot_path, style_json) VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
        (user["id"], name, style.get("primary_color", "#2563eb"),
         style.get("secondary_color", "#64748b"), style.get("accent_color", "#f59e0b"),
         style.get("font_family", "Arial, sans-serif"), filepath, json.dumps(style))
    )
    conn.commit()
    profile_id = cursor.lastrowid
    profile = conn.execute("SELECT * FROM brand_profiles WHERE id = ?", (profile_id,)).fetchone()
    conn.close()

    return {
        "profile": dict(profile),
        "extracted_style": style,
        "message": "Brand style extracted successfully"
    }

@app.get("/api/brand/profiles")
async def get_brand_profiles(user=Depends(get_current_user)):
    """Get all brand profiles for the user"""
    conn = get_db()
    profiles = conn.execute(
        "SELECT * FROM brand_profiles WHERE user_id = ? ORDER BY created_at DESC",
        (user["id"],)
    ).fetchall()
    conn.close()
    return {"profiles": [dict(p) for p in profiles]}

@app.delete("/api/brand/profiles/{profile_id}")
async def delete_brand_profile(profile_id: int, user=Depends(get_current_user)):
    conn = get_db()
    conn.execute("DELETE FROM brand_profiles WHERE id = ? AND user_id = ?", (profile_id, user["id"]))
    conn.commit()
    conn.close()
    return {"message": "Brand profile deleted"}

# ==================== DOCUMENT ENDPOINTS ====================

@app.get("/api/documents")
async def get_user_documents(user=Depends(get_current_user)):
    conn = get_db()
    docs = conn.execute(
        "SELECT id, title, skill_used, prompt, model_used, created_at, updated_at "
        "FROM documents WHERE user_id = ? ORDER BY updated_at DESC",
        (user["id"],)
    ).fetchall()
    conn.close()
    return {"documents": [dict(d) for d in docs]}

@app.get("/api/documents/{doc_id}")
async def get_document(doc_id: int, user=Depends(get_current_user)):
    conn = get_db()
    doc = conn.execute("SELECT * FROM documents WHERE id = ? AND user_id = ?", (doc_id, user["id"])).fetchone()
    conn.close()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    return dict(doc)

@app.get("/api/documents/{doc_id}/pdf")
async def get_document_pdf(doc_id: int, user=Depends(get_current_user)):
    """Download a document as PDF"""
    conn = get_db()
    doc = conn.execute("SELECT * FROM documents WHERE id = ? AND user_id = ?", (doc_id, user["id"])).fetchone()
    conn.close()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")

    # Get brand style if linked
    brand_style = json.loads(doc["brand_style"]) if doc["brand_style"] else None

    html = doc["html_content"] or markdown_to_html(doc["content"])
    loop = asyncio.get_event_loop()
    pdf_bytes = await loop.run_in_executor(_executor, generate_pdf_bytes, html, brand_style)

    return Response(
        content=pdf_bytes,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{doc["title"]}.pdf"'}
    )

@app.put("/api/documents/{doc_id}")
async def update_document(doc_id: int, update: DocumentUpdate, user=Depends(get_current_user)):
    conn = get_db()
    doc = conn.execute("SELECT id FROM documents WHERE id = ? AND user_id = ?", (doc_id, user["id"])).fetchone()
    if not doc:
        conn.close()
        raise HTTPException(status_code=404, detail="Document not found")

    updates, values = [], []
    if update.title is not None:
        updates.append("title = ?")
        values.append(update.title)
    if update.content is not None:
        updates.append("content = ?")
        values.append(update.content)
        updates.append("html_content = ?")
        values.append(markdown_to_html(update.content))

    if updates:
        updates.append("updated_at = datetime('now')")
        values.extend([doc_id, user["id"]])
        conn.execute(f"UPDATE documents SET {', '.join(updates)} WHERE id = ? AND user_id = ?", values)
        conn.commit()

        # Save version snapshot
        version_num = conn.execute("SELECT COUNT(*) as c FROM document_versions WHERE document_id = ?", (doc_id,)).fetchone()["c"] + 1
        conn.execute("INSERT INTO document_versions (document_id, content, version_number, created_by) VALUES (?, ?, ?, ?)",
            (doc_id, update.content if update.content is not None else "", version_num, user["id"]))
        conn.commit()

    updated = conn.execute("SELECT * FROM documents WHERE id = ?", (doc_id,)).fetchone()
    conn.close()
    return dict(updated)

@app.delete("/api/documents/{doc_id}")
async def delete_document(doc_id: int, user=Depends(get_current_user)):
    conn = get_db()
    doc = conn.execute("SELECT id FROM documents WHERE id = ? AND user_id = ?", (doc_id, user["id"])).fetchone()
    if not doc:
        conn.close()
        raise HTTPException(status_code=404, detail="Document not found")
    conn.execute("DELETE FROM documents WHERE id = ?", (doc_id,))
    conn.commit()
    conn.close()
    return {"message": "Document deleted"}

# ==================== VERSION HISTORY ENDPOINTS ====================

@app.get("/api/documents/{doc_id}/versions")
async def get_document_versions(doc_id: int, user=Depends(get_current_user)):
    conn = get_db()
    doc = conn.execute("SELECT user_id FROM documents WHERE id = ?", (doc_id,)).fetchone()
    if not doc or doc["user_id"] != user["id"]:
        conn.close()
        raise HTTPException(status_code=404, detail="Document not found")
    versions = conn.execute(
        "SELECT id, version_number, created_at, LENGTH(content) as content_length FROM document_versions WHERE document_id = ? ORDER BY version_number DESC",
        (doc_id,)
    ).fetchall()
    conn.close()
    return {"versions": [dict(v) for v in versions]}

@app.get("/api/documents/{doc_id}/versions/{version_id}")
async def get_document_version(doc_id: int, version_id: int, user=Depends(get_current_user)):
    conn = get_db()
    doc = conn.execute("SELECT user_id FROM documents WHERE id = ?", (doc_id,)).fetchone()
    if not doc or doc["user_id"] != user["id"]:
        conn.close()
        raise HTTPException(status_code=404, detail="Document not found")
    version = conn.execute("SELECT * FROM document_versions WHERE id = ? AND document_id = ?", (version_id, doc_id)).fetchone()
    conn.close()
    if not version:
        raise HTTPException(status_code=404, detail="Version not found")
    return dict(version)

@app.post("/api/documents/{doc_id}/versions/{version_id}/restore")
async def restore_document_version(doc_id: int, version_id: int, user=Depends(get_current_user)):
    conn = get_db()
    doc = conn.execute("SELECT user_id FROM documents WHERE id = ?", (doc_id,)).fetchone()
    if not doc or doc["user_id"] != user["id"]:
        conn.close()
        raise HTTPException(status_code=404, detail="Document not found")
    version = conn.execute("SELECT content FROM document_versions WHERE id = ? AND document_id = ?", (version_id, doc_id)).fetchone()
    if not version:
        conn.close()
        raise HTTPException(status_code=404, detail="Version not found")
    # Save current as new version before restoring
    current = conn.execute("SELECT content FROM documents WHERE id = ?", (doc_id,)).fetchone()
    ver_num = conn.execute("SELECT COUNT(*) as c FROM document_versions WHERE document_id = ?", (doc_id,)).fetchone()["c"] + 1
    conn.execute("INSERT INTO document_versions (document_id, content, version_number, created_by) VALUES (?, ?, ?, ?)",
        (doc_id, current["content"], ver_num, user["id"]))
    # Restore
    conn.execute("UPDATE documents SET content = ?, updated_at = datetime('now') WHERE id = ?", (version["content"], doc_id))
    conn.commit()
    conn.close()
    return {"message": "Version restored", "version_number": ver_num}

# ==================== WEB PUBLISHING ENDPOINTS ====================

@app.post("/api/documents/{doc_id}/publish")
async def publish_document(doc_id: int, request_data: dict = Body({}), user=Depends(get_current_user)):
    conn = get_db()
    doc = conn.execute("SELECT user_id, title FROM documents WHERE id = ?", (doc_id,)).fetchone()
    if not doc or doc["user_id"] != user["id"]:
        conn.close()
        raise HTTPException(status_code=404, detail="Document not found")
    import secrets as sec
    slug = request_data.get("slug") or sec.token_urlsafe(8)
    password = request_data.get("password")
    conn.execute("UPDATE documents SET is_published = 1, publish_slug = ?, publish_password = ? WHERE id = ?",
        (slug, password, doc_id))
    conn.commit()
    conn.close()
    return {"message": "Document published", "slug": slug, "url": f"/p/{slug}"}

@app.post("/api/documents/{doc_id}/unpublish")
async def unpublish_document(doc_id: int, user=Depends(get_current_user)):
    conn = get_db()
    doc = conn.execute("SELECT user_id FROM documents WHERE id = ?", (doc_id,)).fetchone()
    if not doc or doc["user_id"] != user["id"]:
        conn.close()
        raise HTTPException(status_code=404, detail="Document not found")
    conn.execute("UPDATE documents SET is_published = 0 WHERE id = ?", (doc_id,))
    conn.commit()
    conn.close()
    return {"message": "Document unpublished"}

# ==================== GENERATION ENDPOINTS ====================

@app.get("/")
async def root():
    # In desktop mode (dist/ exists), serve the frontend
    dist_dir = os.path.join(os.path.dirname(__file__), "..", "dist")
    index_path = os.path.join(dist_dir, "index.html")
    if os.path.isfile(index_path):
        return FileResponse(index_path, media_type="text/html")
    # Fallback: API info for dev mode
    return {
        "message": "Universal Document Creator API",
        "version": "2.0.0",
        "endpoints": {
            "auth": "/api/auth",
            "skills": "/api/skills",
            "generate": "/api/generate",
            "documents": "/api/documents",
            "brand": "/api/brand",
            "models": "/api/models",
            "settings": "/api/settings"
        }
    }

@app.get("/api/skills")
async def get_skills():
    skills = load_all_skills()
    return {"skills": [skill.model_dump() for skill in skills], "count": len(skills)}

@app.get("/api/skills/{skill_name}")
async def get_skill(skill_name: str):
    skill = get_skill_by_name(skill_name)
    if not skill:
        raise HTTPException(status_code=404, detail=f"Skill '{skill_name}' not found")
    return skill.model_dump()

def render_business_card_html(params: dict) -> str:
    """Render a visual business card as styled HTML — xhtml2pdf compatible."""
    name = params.get("full_name", "Your Name")
    title = params.get("job_title", "Your Title")
    company = params.get("company_name", "Company")
    email = params.get("email", "")
    phone = params.get("phone", "")
    website = params.get("website", "")
    tagline = params.get("tagline", "")
    # xhtml2pdf only supports basic CSS — use tables for layout, solid backgrounds, no gradients
    contact_rows = ""
    if email:
        contact_rows += f'<tr><td style="padding:8px 16px;color:#666666;font-size:10pt;width:70px;">EMAIL</td><td style="padding:8px 16px;font-size:11pt;color:#1a1a2e;font-weight:bold;">{email}</td></tr>'
    if phone:
        contact_rows += f'<tr><td style="padding:8px 16px;color:#666666;font-size:10pt;width:70px;">PHONE</td><td style="padding:8px 16px;font-size:11pt;color:#1a1a2e;font-weight:bold;">{phone}</td></tr>'
    if website:
        contact_rows += f'<tr><td style="padding:8px 16px;color:#666666;font-size:10pt;width:70px;">WEB</td><td style="padding:8px 16px;font-size:11pt;color:#1a1a2e;font-weight:bold;">{website}</td></tr>'
    tagline_html = f'<tr><td colspan="2" style="padding:14px 32px 0 32px;"><table width="100%"><tr><td style="border-top:2px solid #0891b2;padding-top:12px;font-size:10pt;color:#0891b2;font-style:italic;">&ldquo;{tagline}&rdquo;</td></tr></table></td></tr>' if tagline else ''

    return f'''<table width="100%" cellpadding="0" cellspacing="0" style="font-family:Helvetica,Arial,sans-serif;">
<tr><td style="padding:20px 0;">

<!-- FRONT SIDE LABEL -->
<table width="100%" style="margin-bottom:8px;"><tr><td style="font-size:9pt;color:#999999;text-transform:uppercase;letter-spacing:2px;">&#9632; Front Side</td></tr></table>

<!-- FRONT OF CARD -->
<table width="100%" cellpadding="0" cellspacing="0" style="background-color:#0f172a;border:1px solid #0f172a;">
<tr><td style="padding:24px 32px 6px 32px;font-size:9pt;color:#0891b2;letter-spacing:2px;text-transform:uppercase;font-weight:bold;">{company}</td></tr>
<tr><td style="padding:4px 32px 4px 32px;font-size:22pt;color:#ffffff;font-weight:bold;">{name}</td></tr>
<tr><td style="padding:2px 32px 16px 32px;font-size:11pt;color:#94a3b8;text-transform:uppercase;letter-spacing:1px;">{title}</td></tr>
{tagline_html}
<tr><td style="padding:0 32px 4px 32px;"><table width="100%"><tr><td style="border-top:1px solid #1e3a5f;">&nbsp;</td></tr></table></td></tr>
<tr><td style="padding:0 32px 20px 32px;font-size:8pt;color:#475569;">{company} &bull; Innovation &bull; Excellence</td></tr>
</table>

<!-- SPACING -->
<table width="100%"><tr><td style="padding:10px 0;">&nbsp;</td></tr></table>

<!-- BACK SIDE LABEL -->
<table width="100%" style="margin-bottom:8px;"><tr><td style="font-size:9pt;color:#999999;text-transform:uppercase;letter-spacing:2px;">&#9632; Back Side</td></tr></table>

<!-- BACK OF CARD -->
<table width="100%" cellpadding="0" cellspacing="0" style="background-color:#f8fafc;border:1px solid #e2e8f0;">
<tr><td style="padding:20px 32px 8px 32px;">
  <table width="100%"><tr><td style="font-size:9pt;color:#0891b2;letter-spacing:2px;text-transform:uppercase;font-weight:bold;border-bottom:2px solid #0891b2;padding-bottom:10px;">{company}</td></tr></table>
</td></tr>
<tr><td style="padding:8px 16px 20px 16px;">
  <table width="100%" cellpadding="0" cellspacing="0" style="border:1px solid #e2e8f0;">
    {contact_rows}
  </table>
</td></tr>
</table>

<!-- FOOTER -->
<table width="100%" style="margin-top:6px;"><tr><td style="text-align:center;font-size:7pt;color:#cccccc;padding-top:4px;">Generated by Universal Document Creator</td></tr></table>

</td></tr></table>'''


def render_receipt_html(params: dict, items_text: str = "") -> str:
    """Render a professional receipt as styled HTML — no LLM needed for layout."""
    biz = params.get("business_name", "Business")
    customer = params.get("customer_name", "Customer")
    payment = params.get("payment_method", "Cash")
    tax_rate = params.get("tax_rate", "0")
    receipt_no = params.get("receipt_number", f"R-{datetime.utcnow().strftime('%Y%m%d%H%M')}")
    date_str = datetime.utcnow().strftime("%B %d, %Y  %I:%M %p")

    # Parse items from the prompt text (look for price patterns)
    # This is best-effort — the LLM will generate the content, we just style it
    return None  # Let LLM handle receipt content, we just fix the output format


@app.post("/api/generate/stream")
async def generate_document_stream(request: DocumentRequest, user=Depends(get_optional_user)):
    """Stream document generation via SSE — real-time token-by-token output"""
    import json as json_module

    # Check rate limit
    conn = get_db()
    try:
        check_generation_limit(user, conn)
    finally:
        conn.close()

    # Build the prompt (same logic as generate_document)
    prompt = request.prompt
    skill_used = None

    # Apply skill template if selected
    if request.skill_name:
        skill = get_skill_by_name(request.skill_name)
        if skill:
            skill_prompt = apply_skill_template(skill, request.parameters)
            prompt = f"{skill_prompt}\n\nAdditional context: {request.prompt}"
            skill_used = skill.name

    # Get brand style if specified
    brand_style_instruction = ""
    if request.brand_profile_id and user:
        conn = get_db()
        profile = conn.execute(
            "SELECT * FROM brand_profiles WHERE id = ? AND user_id = ?",
            (request.brand_profile_id, user["id"])
        ).fetchone()
        conn.close()
        if profile and profile["style_json"]:
            brand_style = json.loads(profile["style_json"])
            brand_style_instruction = f"\n\nIMPORTANT - Style this document to match the following brand:\n- Primary color: {brand_style.get('primary_color')}\n- Secondary color: {brand_style.get('secondary_color')}\n- Style: {brand_style.get('style_description', 'professional')}\n- Use appropriate HTML formatting with inline styles matching these brand colors."

    # Format instruction
    format_instruction = ""
    if request.output_format == "html":
        format_instruction = "\n\nGenerate this as clean, well-structured HTML with appropriate tags (h1, h2, p, ul, table, etc.). Include inline CSS for styling."
    elif request.output_format == "pdf":
        format_instruction = "\n\nGenerate this as clean HTML suitable for PDF conversion. Use proper heading hierarchy, tables, lists, and structure. Include inline CSS for professional styling."
    else:
        format_instruction = "\n\nGenerate this as well-formatted Markdown with proper headings, lists, tables, and emphasis."

    full_prompt = prompt + format_instruction + brand_style_instruction

    if request.language:
        full_prompt = f"{full_prompt}\n\nIMPORTANT: Write the entire document in {request.language}."

    # Determine provider
    model_id = request.model
    provider = None
    for m in AVAILABLE_MODELS:
        if m["id"] == model_id:
            provider = m["provider"]
            break
    if not provider:
        provider = "google"

    async def event_generator():
        try:
            if provider == "google":
                # Gemini streaming via thread + asyncio.Queue
                api_key = GEMINI_API_KEY
                if not api_key:
                    conn = get_db()
                    row = conn.execute("SELECT value FROM settings WHERE key = 'gemini_api_key'").fetchone()
                    conn.close()
                    if row:
                        api_key = row["value"]
                if not api_key:
                    yield {"event": "error", "data": json_module.dumps({"error": "Gemini API key not configured. Go to Settings to add it."})}
                    return

                queue = asyncio.Queue()

                def _stream_gemini():
                    try:
                        client = genai.Client(api_key=api_key)
                        for chunk in client.models.generate_content_stream(
                            model=model_id,
                            contents=full_prompt,
                            config=genai_types.GenerateContentConfig(
                                temperature=request.temperature,
                                max_output_tokens=request.max_tokens,
                            ),
                        ):
                            if chunk.text:
                                queue.put_nowait(chunk.text)
                        queue.put_nowait(None)  # sentinel
                    except Exception as e:
                        queue.put_nowait(Exception(str(e)))

                loop = asyncio.get_event_loop()
                loop.run_in_executor(_stream_executor, _stream_gemini)

                full_text = ""
                while True:
                    text = await queue.get()
                    if text is None:
                        break
                    if isinstance(text, Exception):
                        yield {"event": "error", "data": json_module.dumps({"error": str(text)})}
                        return
                    full_text += text
                    yield {"event": "token", "data": json_module.dumps({"text": text})}

                yield {"event": "done", "data": json_module.dumps({"full_text": full_text, "model_used": model_id, "skill_used": skill_used})}

            elif provider == "anthropic":
                # Anthropic native SSE streaming
                api_key = ANTHROPIC_API_KEY
                if not api_key:
                    conn = get_db()
                    row = conn.execute("SELECT value FROM settings WHERE key = 'anthropic_api_key'").fetchone()
                    conn.close()
                    if row:
                        api_key = row["value"]
                if not api_key:
                    yield {"event": "error", "data": json_module.dumps({"error": "Anthropic API key not configured. Go to Settings to add it."})}
                    return

                import httpx
                headers = {
                    "x-api-key": api_key,
                    "anthropic-version": "2023-06-01",
                    "content-type": "application/json",
                }
                payload = {
                    "model": model_id,
                    "max_tokens": request.max_tokens,
                    "temperature": request.temperature,
                    "stream": True,
                    "messages": [{"role": "user", "content": full_prompt}],
                }
                full_text = ""
                async with httpx.AsyncClient(timeout=120) as http_client:
                    async with http_client.stream("POST", "https://api.anthropic.com/v1/messages", headers=headers, json=payload) as resp:
                        async for line in resp.aiter_lines():
                            if line.startswith("data: "):
                                try:
                                    data = json_module.loads(line[6:])
                                except json_module.JSONDecodeError:
                                    continue
                                if data.get("type") == "content_block_delta":
                                    text = data["delta"].get("text", "")
                                    if text:
                                        full_text += text
                                        yield {"event": "token", "data": json_module.dumps({"text": text})}

                yield {"event": "done", "data": json_module.dumps({"full_text": full_text, "model_used": model_id, "skill_used": skill_used})}

            else:
                # Ollama — generate full response then stream in word chunks
                result = await generate_with_ollama(full_prompt, model_id, request.temperature)
                words = result.split(" ")
                full_text = ""
                for i in range(0, len(words), 3):
                    chunk = " ".join(words[i:i + 3])
                    if i + 3 < len(words):
                        chunk += " "
                    full_text += chunk
                    yield {"event": "token", "data": json_module.dumps({"text": chunk})}
                    await asyncio.sleep(0.02)

                yield {"event": "done", "data": json_module.dumps({"full_text": full_text, "model_used": model_id, "skill_used": skill_used})}

            # Increment generation count for free users after successful generation
            if user and user.get("plan") == "free" and not user.get("is_admin"):
                inc_conn = get_db()
                inc_conn.execute(
                    "UPDATE users SET generation_count = COALESCE(generation_count, 0) + 1 WHERE id = ?",
                    (user["id"],)
                )
                inc_conn.commit()
                inc_conn.close()

        except Exception as e:
            yield {"event": "error", "data": json_module.dumps({"error": str(e)})}

    return EventSourceResponse(event_generator())


@app.post("/api/generate", response_model=DocumentResponse)
async def generate_document(request: DocumentRequest, user=Depends(get_optional_user)):
    """Generate a document with selected AI model"""
    # Check rate limit
    conn = get_db()
    try:
        check_generation_limit(user, conn)
    finally:
        conn.close()

    prompt = request.prompt
    skill_used = None

    # Business Card — use dedicated visual template, skip LLM for layout
    if request.skill_name and request.skill_name.lower() == "business card" and request.parameters:
        card_html = render_business_card_html(request.parameters)
        skill_used = "Business Card"
        content = f"# Business Card — {request.parameters.get('full_name', '')}\n\n{request.parameters.get('job_title', '')} at {request.parameters.get('company_name', '')}"
        html_content = card_html

        # Save + return early
        document_id = None
        if user:
            conn = get_db()
            title = request.title or f"Business Card - {request.parameters.get('full_name', '')}"
            cursor = conn.execute(
                "INSERT INTO documents (user_id, title, content, html_content, skill_used, prompt, model_used, brand_style) VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
                (user["id"], title, content, html_content, skill_used, request.prompt, request.model, None)
            )
            conn.commit()
            document_id = cursor.lastrowid
            # Save initial version
            conn.execute("INSERT INTO document_versions (document_id, content, version_number, created_by) VALUES (?, ?, 1, ?)",
                (document_id, content, user["id"] if user else None))
            conn.commit()
            conn.close()

        timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        try:
            pdf_bytes = generate_pdf_bytes(html_content, None, include_footer=False)
            with open(os.path.join(OUTPUT_DIR, f"Business_Card_{timestamp}.pdf"), "wb") as f:
                f.write(pdf_bytes)
        except Exception:
            pass

        return DocumentResponse(
            content=content, html_content=html_content, skill_used=skill_used,
            model_used="template", export_formats=["pdf", "markdown", "text"],
            generated_at=datetime.utcnow().isoformat(), document_id=document_id
        )

    # Apply skill template if selected
    if request.skill_name:
        skill = get_skill_by_name(request.skill_name)
        if skill:
            skill_prompt = apply_skill_template(skill, request.parameters)
            prompt = f"{skill_prompt}\n\nAdditional context: {request.prompt}"
            skill_used = skill.name
        else:
            raise HTTPException(status_code=404, detail=f"Skill '{request.skill_name}' not found")

    # Get brand style if specified
    brand_style = None
    brand_style_instruction = ""
    if request.brand_profile_id and user:
        conn = get_db()
        profile = conn.execute(
            "SELECT * FROM brand_profiles WHERE id = ? AND user_id = ?",
            (request.brand_profile_id, user["id"])
        ).fetchone()
        conn.close()
        if profile and profile["style_json"]:
            brand_style = json.loads(profile["style_json"])
            brand_style_instruction = f"\n\nIMPORTANT - Style this document to match the following brand:\n- Primary color: {brand_style.get('primary_color')}\n- Secondary color: {brand_style.get('secondary_color')}\n- Style: {brand_style.get('style_description', 'professional')}\n- Use appropriate HTML formatting with inline styles matching these brand colors."

    # Adjust prompt for output format
    format_instruction = ""
    if request.output_format == "html":
        format_instruction = "\n\nGenerate this as clean, well-structured HTML with appropriate tags (h1, h2, p, ul, table, etc.). Include inline CSS for styling."
    elif request.output_format == "pdf":
        format_instruction = "\n\nGenerate this as clean HTML suitable for PDF conversion. Use proper heading hierarchy, tables, lists, and structure. Include inline CSS for professional styling."
    else:
        format_instruction = "\n\nGenerate this as well-formatted Markdown with proper headings, lists, tables, and emphasis."

    full_prompt = prompt + format_instruction + brand_style_instruction

    if request.language:
        full_prompt = f"{full_prompt}\n\nIMPORTANT: Write the entire document in {request.language}."

    # Generate content
    content = await generate_content(
        prompt=full_prompt,
        model_id=request.model,
        temperature=request.temperature,
        max_tokens=request.max_tokens
    )

    # Generate HTML version
    html_content = None
    if request.output_format in ("html", "pdf"):
        html_content = content
    else:
        html_content = markdown_to_html(content)

    # Save document if user is authenticated
    document_id = None
    if user:
        conn = get_db()
        title = request.title or f"Document - {datetime.utcnow().strftime('%b %d, %Y %H:%M')}"
        cursor = conn.execute(
            "INSERT INTO documents (user_id, title, content, html_content, skill_used, prompt, model_used, brand_style) "
            "VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
            (user["id"], title, content, html_content, skill_used, request.prompt,
             request.model, json.dumps(brand_style) if brand_style else None)
        )
        conn.commit()
        document_id = cursor.lastrowid
        # Save initial version
        conn.execute("INSERT INTO document_versions (document_id, content, version_number, created_by) VALUES (?, ?, 1, ?)",
            (document_id, content, user["id"] if user else None))
        conn.commit()
        conn.close()

    # Auto-save to output folder
    timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    safe_skill = (skill_used or "General").replace(" ", "_")
    auto_filename = f"{safe_skill}_{timestamp}"
    try:
        # Save markdown
        md_path = os.path.join(OUTPUT_DIR, f"{auto_filename}.md")
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(content)
        # Save PDF
        pdf_html = html_content or markdown_to_html(content)
        pdf_bytes = generate_pdf_bytes(pdf_html, brand_style)
        pdf_path = os.path.join(OUTPUT_DIR, f"{auto_filename}.pdf")
        with open(pdf_path, "wb") as f:
            f.write(pdf_bytes)
    except Exception:
        pass  # Don't fail the response if auto-save has issues

    # Increment generation count for free users
    if user and user.get("plan") == "free" and not user.get("is_admin"):
        conn = get_db()
        conn.execute(
            "UPDATE users SET generation_count = COALESCE(generation_count, 0) + 1 WHERE id = ?",
            (user["id"],)
        )
        conn.commit()
        conn.close()

    model_id = request.model
    log_audit(user["id"] if user else None, "generate", f"Model: {model_id}, Skill: {skill_used}")

    return DocumentResponse(
        content=content,
        html_content=html_content,
        skill_used=skill_used,
        model_used=request.model,
        generated_at=datetime.utcnow().isoformat(),
        document_id=document_id
    )

@app.post("/api/refine", response_model=DocumentResponse)
async def refine_document(request: RefinementRequest, user=Depends(get_optional_user)):
    refinement_prompt = f"""Please refine the following document based on the feedback provided.

ORIGINAL DOCUMENT:
{request.previous_content}

FEEDBACK FOR IMPROVEMENT:
{request.feedback}

Please provide an improved version that addresses all the feedback while maintaining the original purpose and structure."""

    content = await generate_content(refinement_prompt, request.model)

    return DocumentResponse(
        content=content,
        html_content=markdown_to_html(content),
        skill_used=request.skill_name,
        model_used=request.model,
        generated_at=datetime.utcnow().isoformat()
    )

@app.post("/api/chain")
async def chain_skills(request: SkillChainRequest):
    results = []
    current_content = request.initial_prompt

    for step in request.chain:
        skill_name = step.get("skill_name")
        parameters = step.get("parameters", {})
        skill = get_skill_by_name(skill_name)
        if not skill:
            raise HTTPException(status_code=404, detail=f"Skill '{skill_name}' not found")

        parameters["previous_content"] = current_content
        prompt = apply_skill_template(skill, parameters)
        content = await generate_content(prompt, request.model)
        results.append({"skill": skill_name, "content": content})
        current_content = content

    return {"chain_results": results, "final_output": current_content, "steps": len(results)}

@app.post("/api/skills/upload")
async def upload_skill(file: UploadFile = File(...)):
    if not file.filename.endswith('.md'):
        raise HTTPException(status_code=400, detail="Only .md files are allowed")

    timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    filename = f"{timestamp}_{file.filename}"
    filepath = os.path.join(UPLOADED_SKILLS_DIR, filename)

    try:
        content = await file.read()
        async with aiofiles.open(filepath, 'wb') as f:
            await f.write(content)
        skills = parse_skills_from_markdown(content.decode('utf-8'))
        return {
            "message": "Skill uploaded successfully",
            "filename": filename,
            "skills_added": len(skills),
            "skill_names": [s.name for s in skills]
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error uploading skill: {str(e)}")

# ==================== EXPORT ENDPOINTS ====================

@app.post("/api/export/pdf")
async def export_as_pdf(
    content: str = Form(...),
    filename: Optional[str] = Form(None),
    brand_profile_id: Optional[int] = Form(None),
    user=Depends(get_optional_user)
):
    """Export content as a styled PDF"""
    brand_style = None
    if brand_profile_id and user:
        conn = get_db()
        profile = conn.execute(
            "SELECT style_json FROM brand_profiles WHERE id = ? AND user_id = ?",
            (brand_profile_id, user["id"])
        ).fetchone()
        conn.close()
        if profile and profile["style_json"]:
            brand_style = json.loads(profile["style_json"])

    # If content is already HTML (from a template), use directly; skip extra footer if it has one
    is_raw_html = content.strip().startswith("<")
    html = content if is_raw_html else markdown_to_html(content)
    has_own_footer = "Generated by Universal Document Creator" in content
    loop = asyncio.get_event_loop()
    pdf_bytes = await loop.run_in_executor(_executor, generate_pdf_bytes, html, brand_style, not has_own_footer)

    base_filename = filename or f"document_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}"

    # Auto-save copy to output folder
    save_path = os.path.join(OUTPUT_DIR, f"{base_filename}.pdf")
    with open(save_path, "wb") as f:
        f.write(pdf_bytes)

    log_audit(None, "export_pdf", f"Filename: {base_filename}")
    return Response(
        content=pdf_bytes,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{base_filename}.pdf"'}
    )

@app.post("/api/export/docx")
async def export_docx(
    content: str = Form(...),
    filename: Optional[str] = Form(None),
    user=Depends(get_optional_user)
):
    """Export document as DOCX"""
    from docx import Document as DocxDocument
    from docx.shared import Pt, Inches
    import io

    doc = DocxDocument()
    # Parse markdown-like content
    for line in content.split('\n'):
        line = line.strip()
        if line.startswith('# '):
            doc.add_heading(line[2:], level=1)
        elif line.startswith('## '):
            doc.add_heading(line[3:], level=2)
        elif line.startswith('### '):
            doc.add_heading(line[4:], level=3)
        elif line.startswith('- ') or line.startswith('* '):
            doc.add_paragraph(line[2:], style='List Bullet')
        elif line:
            # Handle bold text **text**
            clean = re.sub(r'\*\*(.*?)\*\*', r'\1', line)
            doc.add_paragraph(clean)

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    base_filename = filename or f"document_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}"

    # Auto-save copy to output folder
    save_path = os.path.join(OUTPUT_DIR, f"{base_filename}.docx")
    with open(save_path, "wb") as f:
        f.write(buffer.getvalue())
    buffer.seek(0)

    log_audit(None, "export_docx", f"Filename: {base_filename}")
    return Response(
        content=buffer.read(),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{base_filename}.docx"'}
    )

@app.post("/api/export/html")
async def export_html(content: str = Form(...), filename: Optional[str] = Form(None)):
    """Export document as styled HTML"""
    html = markdown_to_html(content)
    styled_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{filename or 'Document'}</title>
<style>
body {{ font-family: 'Segoe UI', system-ui, sans-serif; max-width: 800px; margin: 40px auto; padding: 20px; line-height: 1.6; color: #1a1a2e; }}
h1 {{ color: #0f172a; border-bottom: 2px solid #e2e8f0; padding-bottom: 8px; }}
h2 {{ color: #1e293b; }}
h3 {{ color: #334155; }}
table {{ border-collapse: collapse; width: 100%; margin: 16px 0; }}
th, td {{ border: 1px solid #e2e8f0; padding: 8px 12px; text-align: left; }}
th {{ background: #f8fafc; font-weight: 600; }}
code {{ background: #f1f5f9; padding: 2px 6px; border-radius: 4px; font-size: 0.9em; }}
pre {{ background: #f8fafc; padding: 16px; border-radius: 8px; overflow-x: auto; }}
blockquote {{ border-left: 4px solid #e2e8f0; margin: 16px 0; padding: 8px 16px; color: #64748b; }}
</style>
</head>
<body>
{html}
</body>
</html>"""

    base_filename = filename or f"document_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}"
    return Response(
        content=styled_html,
        media_type="text/html",
        headers={"Content-Disposition": f'attachment; filename="{base_filename}.html"'}
    )

@app.post("/api/export/{format}")
async def export_document(format: str, content: str = Form(...), filename: Optional[str] = Form(None)):
    """Export document in text or markdown format"""
    if format not in ["text", "markdown"]:
        raise HTTPException(status_code=400, detail="Format must be text, markdown, or pdf")

    timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    base_filename = filename or f"document_{timestamp}"

    if format == "text":
        ext = "txt"
    else:
        ext = "md"

    # Auto-save copy to output folder
    save_path = os.path.join(OUTPUT_DIR, f"{base_filename}.{ext}")
    with open(save_path, "w", encoding="utf-8") as f:
        f.write(content)

    return {"content": content, "filename": f"{base_filename}.{ext}", "mime_type": f"text/{'plain' if format == 'text' else 'markdown'}", "format": format}

@app.get("/api/models")
async def get_available_models():
    """Get list of available AI models with their capabilities"""
    def _check_models():
        results = []
        for model in AVAILABLE_MODELS:
            model_info = dict(model)
            if model["provider"] == "ollama":
                try:
                    ollama_client.show(model=model["id"])
                    model_info["installed"] = True
                except Exception:
                    model_info["installed"] = False
            else:
                model_info["installed"] = True
            results.append(model_info)
        return results

    loop = asyncio.get_event_loop()
    models_with_status = await loop.run_in_executor(_executor, _check_models)
    return {"models": models_with_status, "default": "gemini-3-flash-preview"}

@app.get("/api/image-models")
async def get_image_models():
    """Get available image generation models"""
    return {"models": IMAGE_MODELS, "default": "gemini-3-pro-image-preview"}


@app.post("/api/generate-image")
async def generate_image_endpoint(
    prompt: str = Form(...),
    model: str = Form("gemini-3-pro-image-preview"),
    aspect_ratio: str = Form("1:1"),
    filename: Optional[str] = Form(None)
):
    """Generate an image from a text prompt"""
    image_bytes = await generate_image(prompt, model, aspect_ratio)

    base_filename = filename or f"image_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}"
    save_path = os.path.join(OUTPUT_DIR, f"{base_filename}.png")
    with open(save_path, "wb") as f:
        f.write(image_bytes)

    # Return base64 + save path
    return {
        "image_base64": base64.b64encode(image_bytes).decode("utf-8"),
        "mime_type": "image/png",
        "saved_to": save_path,
        "filename": f"{base_filename}.png"
    }


# ==================== ADMIN & RATE LIMITS ====================

@app.post("/api/admin/promote")
async def promote_to_admin(request_data: dict = Body(...)):
    """Promote a user to admin. Requires admin secret."""
    admin_secret = os.getenv("ADMIN_SECRET", "universaldoc-admin-2026")

    if request_data.get("secret") != admin_secret:
        raise HTTPException(status_code=403, detail="Invalid admin secret")

    email = request_data.get("email", "").strip()
    if not email:
        raise HTTPException(status_code=400, detail="Email required")

    conn = get_db()
    user = conn.execute("SELECT id FROM users WHERE email = ?", (email,)).fetchone()
    if not user:
        conn.close()
        raise HTTPException(status_code=404, detail="User not found")

    conn.execute("UPDATE users SET is_admin = 1, plan = 'enterprise' WHERE id = ?", (user["id"],))
    conn.commit()
    conn.close()

    return {"message": f"User {email} promoted to admin with enterprise plan"}


@app.get("/api/user/dashboard")
async def get_user_dashboard(user=Depends(get_current_user)):
    """Get dashboard stats for the current user"""
    conn = get_db()
    user_id = user["id"]

    # Document stats
    total_docs = conn.execute("SELECT COUNT(*) as c FROM documents WHERE user_id = ?", (user_id,)).fetchone()["c"]
    recent_docs = conn.execute(
        "SELECT COUNT(*) as c FROM documents WHERE user_id = ? AND created_at > datetime('now', '-7 days')",
        (user_id,)
    ).fetchone()["c"]

    # Generation stats
    row = conn.execute(
        "SELECT plan, is_admin, generation_count, generation_reset, created_at FROM users WHERE id = ?",
        (user_id,)
    ).fetchone()

    # Skill usage breakdown
    skill_stats = conn.execute(
        "SELECT skill_used, COUNT(*) as count FROM documents WHERE user_id = ? AND skill_used IS NOT NULL GROUP BY skill_used ORDER BY count DESC LIMIT 5",
        (user_id,)
    ).fetchall()

    # Model usage breakdown
    model_stats = conn.execute(
        "SELECT model_used, COUNT(*) as count FROM documents WHERE user_id = ? AND model_used IS NOT NULL GROUP BY model_used ORDER BY count DESC LIMIT 5",
        (user_id,)
    ).fetchall()

    conn.close()

    is_unlimited = row["is_admin"] or row["plan"] in ("pro", "enterprise")

    return {
        "plan": row["plan"],
        "is_admin": bool(row["is_admin"]),
        "member_since": row["created_at"],
        "total_documents": total_docs,
        "documents_this_week": recent_docs,
        "generations_used": row["generation_count"] or 0,
        "generations_limit": None if is_unlimited else 10,
        "unlimited": is_unlimited,
        "top_skills": [{"name": r["skill_used"], "count": r["count"]} for r in skill_stats],
        "top_models": [{"name": r["model_used"], "count": r["count"]} for r in model_stats],
    }


@app.get("/api/user/limits")
async def get_user_limits(user=Depends(get_current_user)):
    """Get current user's generation limits"""
    conn = get_db()
    row = conn.execute(
        "SELECT plan, is_admin, generation_count, generation_reset FROM users WHERE id = ?",
        (user["id"],)
    ).fetchone()
    conn.close()

    if not row:
        return {"plan": "free", "is_admin": False, "generations_used": 0, "generations_limit": 10, "unlimited": False}

    is_unlimited = row["is_admin"] or row["plan"] in ("pro", "enterprise")

    return {
        "plan": row["plan"],
        "is_admin": bool(row["is_admin"]),
        "generations_used": row["generation_count"] or 0,
        "generations_limit": None if is_unlimited else 10,
        "unlimited": is_unlimited,
    }


# ==================== TEAM ENDPOINTS ====================

@app.post("/api/teams")
async def create_team(request_data: dict = Body(...), user=Depends(get_current_user)):
    """Create a new team (enterprise plan required)"""
    if user.get("plan") not in ("enterprise",) and not user.get("is_admin"):
        raise HTTPException(status_code=403, detail="Enterprise plan required to create teams")

    name = request_data.get("name", "").strip()
    if not name:
        raise HTTPException(status_code=400, detail="Team name required")

    conn = get_db()
    cursor = conn.execute(
        "INSERT INTO teams (name, owner_id) VALUES (?, ?)",
        (name, user["id"])
    )
    team_id = cursor.lastrowid
    # Add owner as admin member
    conn.execute(
        "INSERT INTO team_members (team_id, user_id, role) VALUES (?, ?, 'admin')",
        (team_id, user["id"])
    )
    conn.commit()
    conn.close()
    return {"id": team_id, "name": name, "message": "Team created"}


@app.get("/api/teams")
async def get_my_teams(user=Depends(get_current_user)):
    """Get teams the user belongs to"""
    conn = get_db()
    teams = conn.execute("""
        SELECT t.id, t.name, t.plan, t.max_members, t.created_at, tm.role,
               (SELECT COUNT(*) FROM team_members WHERE team_id = t.id) as member_count
        FROM teams t
        JOIN team_members tm ON t.id = tm.team_id
        WHERE tm.user_id = ?
    """, (user["id"],)).fetchall()
    conn.close()
    return {"teams": [dict(t) for t in teams]}


@app.get("/api/teams/{team_id}")
async def get_team(team_id: int, user=Depends(get_current_user)):
    """Get team details with members"""
    conn = get_db()
    # Verify user is member
    member = conn.execute(
        "SELECT role FROM team_members WHERE team_id = ? AND user_id = ?",
        (team_id, user["id"])
    ).fetchone()
    if not member:
        conn.close()
        raise HTTPException(status_code=403, detail="Not a member of this team")

    team = conn.execute("SELECT * FROM teams WHERE id = ?", (team_id,)).fetchone()
    members = conn.execute("""
        SELECT u.id, u.name, u.email, tm.role, tm.joined_at
        FROM team_members tm
        JOIN users u ON tm.user_id = u.id
        WHERE tm.team_id = ?
    """, (team_id,)).fetchall()
    conn.close()

    return {
        "team": dict(team),
        "members": [dict(m) for m in members],
        "your_role": member["role"]
    }


@app.post("/api/teams/{team_id}/members")
async def add_team_member(team_id: int, request_data: dict = Body(...), user=Depends(get_current_user)):
    """Add a member to a team (admin only)"""
    conn = get_db()
    # Verify caller is admin
    caller = conn.execute(
        "SELECT role FROM team_members WHERE team_id = ? AND user_id = ?",
        (team_id, user["id"])
    ).fetchone()
    if not caller or caller["role"] != "admin":
        conn.close()
        raise HTTPException(status_code=403, detail="Only team admins can add members")

    email = request_data.get("email", "").strip()
    role = request_data.get("role", "member")
    if role not in ("member", "admin"):
        role = "member"

    target = conn.execute("SELECT id FROM users WHERE email = ?", (email,)).fetchone()
    if not target:
        conn.close()
        raise HTTPException(status_code=404, detail="User not found")

    try:
        conn.execute(
            "INSERT INTO team_members (team_id, user_id, role) VALUES (?, ?, ?)",
            (team_id, target["id"], role)
        )
        conn.commit()
    except Exception:
        conn.close()
        raise HTTPException(status_code=400, detail="User already in team")

    conn.close()
    return {"message": f"Added {email} as {role}"}


@app.delete("/api/teams/{team_id}/members/{member_user_id}")
async def remove_team_member(team_id: int, member_user_id: int, user=Depends(get_current_user)):
    """Remove a member from a team (admin only)"""
    conn = get_db()
    caller = conn.execute(
        "SELECT role FROM team_members WHERE team_id = ? AND user_id = ?",
        (team_id, user["id"])
    ).fetchone()
    if not caller or caller["role"] != "admin":
        conn.close()
        raise HTTPException(status_code=403, detail="Only team admins can remove members")

    conn.execute("DELETE FROM team_members WHERE team_id = ? AND user_id = ?", (team_id, member_user_id))
    conn.commit()
    conn.close()
    return {"message": "Member removed"}


# Shared templates
@app.post("/api/teams/{team_id}/templates")
async def create_shared_template(team_id: int, request_data: dict = Body(...), user=Depends(get_current_user)):
    """Share a template with the team"""
    conn = get_db()
    member = conn.execute(
        "SELECT role FROM team_members WHERE team_id = ? AND user_id = ?",
        (team_id, user["id"])
    ).fetchone()
    if not member:
        conn.close()
        raise HTTPException(status_code=403, detail="Not a member of this team")

    name = request_data.get("name", "").strip()
    content = request_data.get("content", "").strip()
    skill_name = request_data.get("skill_name")

    if not name or not content:
        conn.close()
        raise HTTPException(status_code=400, detail="Name and content required")

    cursor = conn.execute(
        "INSERT INTO shared_templates (team_id, name, content, skill_name, created_by) VALUES (?, ?, ?, ?, ?)",
        (team_id, name, content, skill_name, user["id"])
    )
    conn.commit()
    conn.close()
    return {"id": cursor.lastrowid, "message": "Template shared"}


@app.get("/api/teams/{team_id}/templates")
async def get_shared_templates(team_id: int, user=Depends(get_current_user)):
    """Get shared templates for a team"""
    conn = get_db()
    member = conn.execute(
        "SELECT role FROM team_members WHERE team_id = ? AND user_id = ?",
        (team_id, user["id"])
    ).fetchone()
    if not member:
        conn.close()
        raise HTTPException(status_code=403, detail="Not a member of this team")

    templates = conn.execute("""
        SELECT st.*, u.name as creator_name
        FROM shared_templates st
        JOIN users u ON st.created_by = u.id
        WHERE st.team_id = ?
        ORDER BY st.created_at DESC
    """, (team_id,)).fetchall()
    conn.close()
    return {"templates": [dict(t) for t in templates]}


# ==================== AUDIT LOG ENDPOINTS ====================

@app.get("/api/admin/audit-logs")
async def get_audit_logs(
    user=Depends(get_current_user),
    limit: int = 50,
    offset: int = 0,
    action: Optional[str] = None
):
    """Get audit logs (admin only)"""
    if not user.get("is_admin"):
        raise HTTPException(status_code=403, detail="Admin access required")

    conn = get_db()
    if action:
        logs = conn.execute("""
            SELECT al.*, u.name as user_name, u.email as user_email
            FROM audit_logs al
            LEFT JOIN users u ON al.user_id = u.id
            WHERE al.action = ?
            ORDER BY al.created_at DESC
            LIMIT ? OFFSET ?
        """, (action, limit, offset)).fetchall()
        total = conn.execute("SELECT COUNT(*) as c FROM audit_logs WHERE action = ?", (action,)).fetchone()["c"]
    else:
        logs = conn.execute("""
            SELECT al.*, u.name as user_name, u.email as user_email
            FROM audit_logs al
            LEFT JOIN users u ON al.user_id = u.id
            ORDER BY al.created_at DESC
            LIMIT ? OFFSET ?
        """, (limit, offset)).fetchall()
        total = conn.execute("SELECT COUNT(*) as c FROM audit_logs").fetchone()["c"]

    conn.close()
    return {
        "logs": [dict(l) for l in logs],
        "total": total,
        "limit": limit,
        "offset": offset,
    }


# ==================== BRANDING ENDPOINTS ====================

@app.post("/api/admin/branding")
async def update_branding(request_data: dict = Body(...), user=Depends(get_current_user)):
    """Update site branding (admin only)"""
    if not user.get("is_admin"):
        raise HTTPException(status_code=403, detail="Admin access required")

    conn = get_db()
    branding_keys = ["app_name", "app_tagline", "primary_color", "logo_url", "support_email", "custom_footer"]
    for key in branding_keys:
        if key in request_data:
            conn.execute(
                "INSERT INTO settings (key, value) VALUES (?, ?) ON CONFLICT(key) DO UPDATE SET value = ?",
                (f"branding_{key}", str(request_data[key]), str(request_data[key]))
            )
    conn.commit()
    conn.close()
    log_audit(user["id"], "update_branding", f"Updated: {list(request_data.keys())}")
    return {"message": "Branding updated"}


@app.get("/api/branding")
async def get_branding():
    """Get site branding configuration (public)"""
    conn = get_db()
    rows = conn.execute("SELECT key, value FROM settings WHERE key LIKE 'branding_%'").fetchall()
    conn.close()

    branding = {}
    for row in rows:
        key = row["key"].replace("branding_", "")
        branding[key] = row["value"]

    # Defaults
    defaults = {
        "app_name": "Universal Document Creator",
        "app_tagline": "AI-Powered Document Generation",
        "primary_color": "#ea580c",
        "logo_url": "",
        "support_email": "support@universaldoc.app",
        "custom_footer": "",
    }

    for k, v in defaults.items():
        if k not in branding:
            branding[k] = v

    return branding


# ==================== TEMPLATE MARKETPLACE (#11) ====================

@app.post("/api/marketplace/publish")
async def publish_template(request_data: dict = Body(...), user=Depends(get_current_user)):
    name = request_data.get("name", "").strip()
    description = request_data.get("description", "").strip()
    category = request_data.get("category", "general")
    skill_content = request_data.get("skill_content", "").strip()
    tags = request_data.get("tags", "")
    if not name or not skill_content:
        raise HTTPException(status_code=400, detail="Name and skill content required")
    conn = get_db()
    cursor = conn.execute(
        "INSERT INTO marketplace_templates (author_id, name, description, category, skill_content, tags) VALUES (?, ?, ?, ?, ?, ?)",
        (user["id"], name, description, category, skill_content, tags if isinstance(tags, str) else json.dumps(tags))
    )
    conn.commit()
    conn.close()
    return {"id": cursor.lastrowid, "message": "Template published"}

@app.get("/api/marketplace")
async def browse_marketplace(category: Optional[str] = None, search: Optional[str] = None, sort: str = "popular"):
    conn = get_db()
    query = """SELECT mt.*, u.name as author_name,
        CASE WHEN mt.rating_count > 0 THEN ROUND(CAST(mt.rating_sum AS FLOAT) / mt.rating_count, 1) ELSE 0 END as avg_rating
        FROM marketplace_templates mt JOIN users u ON mt.author_id = u.id WHERE mt.is_approved = 1"""
    params = []
    if category:
        query += " AND mt.category = ?"
        params.append(category)
    if search:
        query += " AND (mt.name LIKE ? OR mt.description LIKE ? OR mt.tags LIKE ?)"
        params.extend([f"%{search}%"] * 3)
    if sort == "popular":
        query += " ORDER BY mt.downloads DESC"
    elif sort == "rating":
        query += " ORDER BY avg_rating DESC"
    elif sort == "newest":
        query += " ORDER BY mt.created_at DESC"
    else:
        query += " ORDER BY mt.downloads DESC"
    query += " LIMIT 50"
    templates = conn.execute(query, params).fetchall()
    conn.close()
    return {"templates": [dict(t) for t in templates]}

@app.get("/api/marketplace/{template_id}")
async def get_marketplace_template(template_id: int):
    conn = get_db()
    t = conn.execute(
        "SELECT mt.*, u.name as author_name FROM marketplace_templates mt JOIN users u ON mt.author_id = u.id WHERE mt.id = ?",
        (template_id,)
    ).fetchone()
    if not t:
        conn.close()
        raise HTTPException(status_code=404, detail="Template not found")
    reviews = conn.execute(
        "SELECT tr.*, u.name as reviewer_name FROM template_reviews tr JOIN users u ON tr.user_id = u.id WHERE tr.template_id = ? ORDER BY tr.created_at DESC LIMIT 10",
        (template_id,)
    ).fetchall()
    conn.close()
    return {"template": dict(t), "reviews": [dict(r) for r in reviews]}

@app.post("/api/marketplace/{template_id}/install")
async def install_marketplace_template(template_id: int, user=Depends(get_current_user)):
    conn = get_db()
    t = conn.execute("SELECT skill_content FROM marketplace_templates WHERE id = ?", (template_id,)).fetchone()
    if not t:
        conn.close()
        raise HTTPException(status_code=404, detail="Template not found")
    conn.execute("UPDATE marketplace_templates SET downloads = downloads + 1 WHERE id = ?", (template_id,))
    conn.commit()
    conn.close()
    return {"skill_content": t["skill_content"], "message": "Template installed"}

@app.post("/api/marketplace/{template_id}/review")
async def review_template(template_id: int, request_data: dict = Body(...), user=Depends(get_current_user)):
    rating = request_data.get("rating", 0)
    review_text = request_data.get("review", "")
    if not 1 <= rating <= 5:
        raise HTTPException(status_code=400, detail="Rating must be 1-5")
    conn = get_db()
    try:
        conn.execute(
            "INSERT INTO template_reviews (template_id, user_id, rating, review) VALUES (?, ?, ?, ?)",
            (template_id, user["id"], rating, review_text)
        )
        conn.execute(
            "UPDATE marketplace_templates SET rating_sum = rating_sum + ?, rating_count = rating_count + 1 WHERE id = ?",
            (rating, template_id)
        )
        conn.commit()
    except Exception:
        conn.close()
        raise HTTPException(status_code=400, detail="You already reviewed this template")
    conn.close()
    return {"message": "Review submitted"}

# ==================== IN-DOCUMENT PAYMENTS (#15) ====================

@app.post("/api/documents/{doc_id}/payment-link")
async def create_document_payment_link(doc_id: int, request_data: dict = Body(...), user=Depends(get_current_user)):
    """Create a Stripe payment link embedded in a document"""
    amount = request_data.get("amount")  # in cents
    description = request_data.get("description", "Document payment")
    currency = request_data.get("currency", "usd")

    if not amount or amount < 50:
        raise HTTPException(status_code=400, detail="Amount must be at least $0.50")

    api_key = STRIPE_SECRET_KEY
    if not api_key:
        conn = get_db()
        row = conn.execute("SELECT value FROM settings WHERE key = 'stripe_secret_key'").fetchone()
        conn.close()
        if row:
            api_key = row["value"]
    if not api_key:
        raise HTTPException(status_code=400, detail="Stripe not configured")

    stripe.api_key = api_key

    try:
        price = stripe.Price.create(
            unit_amount=amount,
            currency=currency,
            product_data={"name": description},
        )
        payment_link = stripe.PaymentLink.create(
            line_items=[{"price": price.id, "quantity": 1}],
            metadata={"document_id": str(doc_id), "user_id": str(user["id"])}
        )
        return {"payment_url": payment_link.url, "amount": amount, "currency": currency}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Stripe error: {str(e)}")

# ==================== AI WORKFLOWS (#16) ====================

@app.get("/api/workflows/templates")
async def get_workflow_templates():
    """Get pre-built workflow templates"""
    return {"workflows": [
        {
            "id": "blog_pipeline",
            "name": "Blog Post Pipeline",
            "description": "Research \u2192 Outline \u2192 Draft \u2192 Edit \u2192 SEO optimize",
            "steps": [
                {"skill_name": None, "prompt_template": "Research the topic: {{topic}} and list 5 key points"},
                {"skill_name": None, "prompt_template": "Create a detailed outline based on these points:\n{{previous_output}}"},
                {"skill_name": None, "prompt_template": "Write a full blog post from this outline:\n{{previous_output}}"},
                {"skill_name": None, "prompt_template": "Edit and improve this blog post for clarity and engagement:\n{{previous_output}}"},
            ]
        },
        {
            "id": "proposal_workflow",
            "name": "Business Proposal Workflow",
            "description": "Company research \u2192 Problem statement \u2192 Solution \u2192 Pricing \u2192 Full proposal",
            "steps": [
                {"skill_name": None, "prompt_template": "Research company {{company_name}} and summarize their industry, challenges, and needs"},
                {"skill_name": "Business Proposal", "prompt_template": "Create a business proposal for {{company_name}} based on this research:\n{{previous_output}}"},
            ]
        },
        {
            "id": "content_repurpose",
            "name": "Content Repurposing",
            "description": "Take one piece of content and create multiple formats",
            "steps": [
                {"skill_name": None, "prompt_template": "Summarize the key points from this content:\n{{source_content}}"},
                {"skill_name": "Email Template", "prompt_template": "Create an email newsletter based on:\n{{previous_output}}"},
                {"skill_name": "Marketing Copy", "prompt_template": "Create 3 social media posts (Twitter, LinkedIn, Instagram) from:\n{{previous_output}}"},
            ]
        },
        {
            "id": "legal_review",
            "name": "Legal Document Review",
            "description": "Analyze \u2192 Identify issues \u2192 Suggest improvements \u2192 Rewrite",
            "steps": [
                {"skill_name": None, "prompt_template": "Analyze this legal document and identify the key terms, obligations, and potential issues:\n{{document_text}}"},
                {"skill_name": None, "prompt_template": "Based on this analysis, suggest specific improvements and flag any risky clauses:\n{{previous_output}}"},
                {"skill_name": "Legal Contract", "prompt_template": "Rewrite the document incorporating the suggested improvements:\n{{previous_output}}"},
            ]
        },
    ]}

@app.post("/api/workflows/execute")
async def execute_workflow(request_data: dict = Body(...), user=Depends(get_optional_user)):
    """Execute a multi-step workflow"""
    steps = request_data.get("steps", [])
    variables = request_data.get("variables", {})
    model = request_data.get("model", "gemini-3-flash-preview")

    if not steps:
        raise HTTPException(status_code=400, detail="No steps provided")

    results = []
    previous_output = ""

    for i, step in enumerate(steps):
        prompt = step.get("prompt_template", "")
        # Replace variables
        for key, value in variables.items():
            prompt = prompt.replace("{{" + key + "}}", str(value))
        prompt = prompt.replace("{{previous_output}}", previous_output)

        try:
            result = await generate_content(prompt, model_id=model, temperature=0.7, max_tokens=2000)
            previous_output = result
            results.append({"step": i + 1, "output": result, "status": "success"})
        except Exception as e:
            results.append({"step": i + 1, "output": str(e), "status": "error"})
            break

    return {
        "results": results,
        "final_output": previous_output,
        "steps_completed": len(results),
        "total_steps": len(steps)
    }

# ==================== SEO ANALYSIS (#8) ====================

@app.post("/api/analyze/seo")
async def analyze_seo(request_data: dict = Body(...)):
    """Analyze content for SEO metrics"""
    content = request_data.get("content", "")
    target_keyword = request_data.get("keyword", "")

    if not content:
        raise HTTPException(status_code=400, detail="Content required")

    import textstat

    # Basic metrics
    word_count = len(content.split())
    sentence_count = textstat.sentence_count(content)

    # Readability scores
    flesch_reading = textstat.flesch_reading_ease(content)
    flesch_grade = textstat.flesch_kincaid_grade(content)
    gunning_fog = textstat.gunning_fog(content)

    # Reading time (avg 200 wpm)
    reading_time_minutes = round(word_count / 200, 1)

    # Keyword analysis
    keyword_data = None
    if target_keyword:
        keyword_lower = target_keyword.lower()
        content_lower = content.lower()
        keyword_count = content_lower.count(keyword_lower)
        keyword_density = round((keyword_count / max(word_count, 1)) * 100, 2)

        # Check keyword in headings
        lines = content.split('\n')
        in_h1 = any(keyword_lower in line.lower() for line in lines if line.startswith('# '))
        in_h2 = any(keyword_lower in line.lower() for line in lines if line.startswith('## '))
        in_first_100 = keyword_lower in content_lower[:500]

        keyword_data = {
            "keyword": target_keyword,
            "count": keyword_count,
            "density_percent": keyword_density,
            "in_title": in_h1,
            "in_subheading": in_h2,
            "in_introduction": in_first_100,
            "recommendation": "Good" if 1 <= keyword_density <= 3 else ("Too low — add more mentions" if keyword_density < 1 else "Too high — reduce to avoid keyword stuffing")
        }

    # Heading structure
    headings = {"h1": 0, "h2": 0, "h3": 0}
    for line in content.split('\n'):
        if line.startswith('### '): headings["h3"] += 1
        elif line.startswith('## '): headings["h2"] += 1
        elif line.startswith('# '): headings["h1"] += 1

    # SEO score (0-100)
    score = 50
    if word_count >= 300: score += 10
    if word_count >= 1000: score += 5
    if headings["h1"] == 1: score += 10
    if headings["h2"] >= 2: score += 5
    if 30 <= flesch_reading <= 70: score += 10  # Good readability
    if keyword_data and 1 <= keyword_data["density_percent"] <= 3: score += 10
    if keyword_data and keyword_data["in_title"]: score += 5
    if keyword_data and keyword_data["in_introduction"]: score += 5
    score = min(score, 100)

    # Readability label
    if flesch_reading >= 60: readability = "Easy to read"
    elif flesch_reading >= 30: readability = "Moderate"
    else: readability = "Difficult"

    return {
        "seo_score": score,
        "word_count": word_count,
        "sentence_count": sentence_count,
        "reading_time_minutes": reading_time_minutes,
        "readability": {
            "flesch_reading_ease": round(flesch_reading, 1),
            "flesch_kincaid_grade": round(flesch_grade, 1),
            "gunning_fog_index": round(gunning_fog, 1),
            "label": readability
        },
        "headings": headings,
        "keyword_analysis": keyword_data,
        "suggestions": []  # Could be expanded with AI-powered suggestions
    }


# ==================== BRAND VOICE TRAINING (#10) ====================

@app.post("/api/voice/analyze")
async def analyze_writing_voice(request_data: dict = Body(...), user=Depends(get_current_user)):
    """Analyze writing samples to extract voice profile"""
    samples = request_data.get("samples", [])
    name = request_data.get("name", "My Voice")

    if not samples or len(samples) < 1:
        raise HTTPException(status_code=400, detail="Provide at least 1 writing sample")

    combined = "\n\n---\n\n".join(samples)

    analysis_prompt = f"""Analyze the following writing samples and extract the author's writing voice profile. Return a JSON object with these exact keys:
- "tone": one of "professional", "casual", "academic", "friendly", "authoritative", "conversational"
- "formality": one of "very formal", "formal", "neutral", "informal", "very informal"
- "vocabulary_level": one of "simple", "moderate", "advanced", "technical"
- "avg_sentence_length": one of "short (under 15 words)", "medium (15-25 words)", "long (over 25 words)"
- "common_phrases": array of 3-5 phrases or patterns the author frequently uses
- "writing_rules": array of 5-8 specific writing rules to follow to match this voice (e.g., "Always use active voice", "Prefer short paragraphs", "Use bullet points for lists")

Writing samples:
{combined}

Return ONLY valid JSON, no markdown formatting."""

    try:
        result = await generate_content(analysis_prompt, model_id="gemini-3-flash-preview", temperature=0.3, max_tokens=1000)
        # Try to parse JSON from result
        import json as json_mod
        # Strip markdown code blocks if present
        clean = result.strip()
        if clean.startswith("```"): clean = clean.split("\n", 1)[1].rsplit("```", 1)[0]
        analysis = json_mod.loads(clean)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to analyze voice: {str(e)}")

    # Save to DB
    conn = get_db()
    cursor = conn.execute(
        """INSERT INTO voice_profiles (user_id, name, tone, formality, vocabulary_level, avg_sentence_length, common_phrases, writing_rules, raw_analysis, sample_count)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
        (user["id"], name,
         analysis.get("tone", ""),
         analysis.get("formality", ""),
         analysis.get("vocabulary_level", ""),
         analysis.get("avg_sentence_length", ""),
         json.dumps(analysis.get("common_phrases", [])),
         json.dumps(analysis.get("writing_rules", [])),
         json.dumps(analysis),
         len(samples))
    )
    conn.commit()
    profile_id = cursor.lastrowid
    conn.close()

    return {"id": profile_id, "name": name, "analysis": analysis}

@app.get("/api/voice/profiles")
async def get_voice_profiles(user=Depends(get_current_user)):
    conn = get_db()
    profiles = conn.execute(
        "SELECT id, name, tone, formality, vocabulary_level, sample_count, created_at FROM voice_profiles WHERE user_id = ? ORDER BY created_at DESC",
        (user["id"],)
    ).fetchall()
    conn.close()
    return {"profiles": [dict(p) for p in profiles]}

@app.get("/api/voice/profiles/{profile_id}")
async def get_voice_profile(profile_id: int, user=Depends(get_current_user)):
    conn = get_db()
    profile = conn.execute("SELECT * FROM voice_profiles WHERE id = ? AND user_id = ?", (profile_id, user["id"])).fetchone()
    conn.close()
    if not profile:
        raise HTTPException(status_code=404, detail="Voice profile not found")
    result = dict(profile)
    for key in ["common_phrases", "writing_rules", "raw_analysis"]:
        if result.get(key):
            try: result[key] = json.loads(result[key])
            except: pass
    return result

@app.delete("/api/voice/profiles/{profile_id}")
async def delete_voice_profile(profile_id: int, user=Depends(get_current_user)):
    conn = get_db()
    conn.execute("DELETE FROM voice_profiles WHERE id = ? AND user_id = ?", (profile_id, user["id"]))
    conn.commit()
    conn.close()
    return {"message": "Voice profile deleted"}


# ==================== MULTI-LANGUAGE (#12) ====================

@app.get("/api/languages")
async def get_supported_languages():
    return {"languages": [
        "English", "Spanish", "French", "German", "Italian", "Portuguese",
        "Dutch", "Russian", "Chinese (Simplified)", "Chinese (Traditional)",
        "Japanese", "Korean", "Arabic", "Hindi", "Turkish", "Polish",
        "Swedish", "Norwegian", "Danish", "Finnish", "Greek", "Czech",
        "Romanian", "Hungarian", "Thai", "Vietnamese", "Indonesian", "Malay"
    ]}


# ==================== BUILT-IN E-SIGNATURES ====================

@app.post("/api/documents/{doc_id}/signature-request")
async def create_signature_request(doc_id: int, request_data: dict = Body(...), user=Depends(get_current_user)):
    """Create a signature request for a document"""
    conn = get_db()
    doc = conn.execute("SELECT user_id, title FROM documents WHERE id = ?", (doc_id,)).fetchone()
    if not doc or doc["user_id"] != user["id"]:
        conn.close()
        raise HTTPException(status_code=404, detail="Document not found")

    import secrets as sec
    sign_token = sec.token_urlsafe(32)
    recipient_email = request_data.get("recipient_email", "").strip()
    recipient_name = request_data.get("recipient_name", "").strip()

    if not recipient_email:
        conn.close()
        raise HTTPException(status_code=400, detail="Recipient email required")

    cursor = conn.execute(
        "INSERT INTO signature_requests (document_id, sender_id, recipient_email, recipient_name, sign_token) VALUES (?, ?, ?, ?, ?)",
        (doc_id, user["id"], recipient_email, recipient_name, sign_token)
    )
    conn.commit()
    conn.close()

    return {"id": cursor.lastrowid, "sign_token": sign_token, "sign_url": f"/sign/{sign_token}", "status": "pending"}

@app.get("/api/documents/{doc_id}/signature-requests")
async def get_signature_requests(doc_id: int, user=Depends(get_current_user)):
    """Get all signature requests for a document"""
    conn = get_db()
    doc = conn.execute("SELECT user_id FROM documents WHERE id = ?", (doc_id,)).fetchone()
    if not doc or doc["user_id"] != user["id"]:
        conn.close()
        raise HTTPException(status_code=404, detail="Document not found")
    requests_list = conn.execute(
        "SELECT id, recipient_email, recipient_name, status, signed_at, created_at FROM signature_requests WHERE document_id = ? ORDER BY created_at DESC",
        (doc_id,)
    ).fetchall()
    conn.close()
    return {"requests": [dict(r) for r in requests_list]}

@app.get("/sign/{sign_token}")
async def view_signature_page(sign_token: str):
    """Public page where recipient can view and sign document"""
    conn = get_db()
    req = conn.execute(
        "SELECT sr.*, d.title, d.content, d.html_content, u.name as sender_name FROM signature_requests sr JOIN documents d ON sr.document_id = d.id JOIN users u ON sr.sender_id = u.id WHERE sr.sign_token = ?",
        (sign_token,)
    ).fetchone()
    conn.close()
    if not req:
        raise HTTPException(status_code=404, detail="Signature request not found")
    if req["status"] == "signed":
        return HTMLResponse(f"""<!DOCTYPE html><html><head><meta charset="UTF-8"><title>Already Signed</title>
<style>body{{font-family:'Segoe UI',sans-serif;max-width:600px;margin:80px auto;text-align:center;color:#1a1a2e}}.badge{{background:#10b981;color:white;padding:8px 24px;border-radius:20px;display:inline-block;margin:20px 0}}</style>
</head><body><h1>{req["title"]}</h1><div class="badge">Signed on {req["signed_at"]}</div><p>This document was signed by {req["recipient_name"] or req["recipient_email"]}.</p></body></html>""")

    html_content = req["html_content"] or markdown_to_html(req["content"])
    return HTMLResponse(f"""<!DOCTYPE html><html><head><meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Sign: {req["title"]}</title>
<style>
body{{font-family:'Segoe UI',sans-serif;max-width:800px;margin:0 auto;padding:20px;color:#1a1a2e}}
.doc-content{{border:1px solid #e2e8f0;padding:24px;border-radius:8px;margin:20px 0;line-height:1.6}}
h1{{color:#0f172a}}
.sign-section{{background:#f8fafc;border:2px dashed #e2e8f0;border-radius:12px;padding:24px;text-align:center;margin:24px 0}}
canvas{{border:2px solid #0f172a;border-radius:8px;cursor:crosshair;background:white;touch-action:none}}
button{{padding:12px 32px;border:none;border-radius:8px;cursor:pointer;font-size:16px;font-weight:600}}
.btn-sign{{background:#ea580c;color:white}}.btn-sign:hover{{background:#c2410c}}
.btn-clear{{background:#e2e8f0;color:#334155;margin-right:12px}}.btn-clear:hover{{background:#cbd5e1}}
.info{{color:#64748b;font-size:14px;margin:8px 0}}
</style></head><body>
<h1>Signature Request</h1>
<p class="info">From: {req["sender_name"]} &bull; To: {req["recipient_name"] or req["recipient_email"]}</p>
<div class="doc-content">{html_content}</div>
<div class="sign-section">
<h3>Your Signature</h3>
<p class="info">Draw your signature below</p>
<canvas id="sigCanvas" width="500" height="200"></canvas><br><br>
<button class="btn-clear" onclick="clearSig()">Clear</button>
<button class="btn-sign" onclick="submitSig()">Sign Document</button>
</div>
<script>
const c=document.getElementById('sigCanvas'),ctx=c.getContext('2d');let drawing=false,lastX=0,lastY=0;
ctx.strokeStyle='#0f172a';ctx.lineWidth=2;ctx.lineCap='round';
c.addEventListener('mousedown',e=>{{drawing=true;[lastX,lastY]=[e.offsetX,e.offsetY]}});
c.addEventListener('mousemove',e=>{{if(!drawing)return;ctx.beginPath();ctx.moveTo(lastX,lastY);ctx.lineTo(e.offsetX,e.offsetY);ctx.stroke();[lastX,lastY]=[e.offsetX,e.offsetY]}});
c.addEventListener('mouseup',()=>drawing=false);c.addEventListener('mouseout',()=>drawing=false);
c.addEventListener('touchstart',e=>{{e.preventDefault();const t=e.touches[0],r=c.getBoundingClientRect();drawing=true;[lastX,lastY]=[t.clientX-r.left,t.clientY-r.top]}});
c.addEventListener('touchmove',e=>{{e.preventDefault();if(!drawing)return;const t=e.touches[0],r=c.getBoundingClientRect();const x=t.clientX-r.left,y=t.clientY-r.top;ctx.beginPath();ctx.moveTo(lastX,lastY);ctx.lineTo(x,y);ctx.stroke();[lastX,lastY]=[x,y]}});
c.addEventListener('touchend',()=>drawing=false);
function clearSig(){{ctx.clearRect(0,0,c.width,c.height)}}
function submitSig(){{
  const data=c.toDataURL('image/png');
  if(ctx.getImageData(0,0,c.width,c.height).data.every((v,i)=>i%4===3?true:v===0)){{alert('Please draw your signature first');return}}
  fetch('/api/signatures/sign',{{method:'POST',headers:{{'Content-Type':'application/json'}},body:JSON.stringify({{token:'{sign_token}',signature_data:data}})}})
  .then(r=>r.json()).then(d=>{{if(d.message){{document.body.innerHTML='<div style="text-align:center;padding:80px 20px"><h1 style="color:#10b981">Signed Successfully</h1><p>Thank you for signing this document.</p></div>'}}else{{alert(d.detail||"Error signing")}}}})
  .catch(()=>alert('Failed to submit signature'))
}}
</script></body></html>""")

@app.post("/api/signatures/sign")
async def submit_signature(request_data: dict = Body(...), request: Request = None):
    """Submit a signature for a document"""
    token = request_data.get("token", "")
    signature_data = request_data.get("signature_data", "")
    if not token or not signature_data:
        raise HTTPException(status_code=400, detail="Token and signature required")

    conn = get_db()
    req = conn.execute("SELECT id, status FROM signature_requests WHERE sign_token = ?", (token,)).fetchone()
    if not req:
        conn.close()
        raise HTTPException(status_code=404, detail="Invalid token")
    if req["status"] == "signed":
        conn.close()
        raise HTTPException(status_code=400, detail="Already signed")

    ip = request.client.host if request and request.client else "unknown"
    conn.execute(
        "UPDATE signature_requests SET status = 'signed', signed_at = datetime('now'), signature_data = ?, ip_address = ? WHERE id = ?",
        (signature_data, ip, req["id"])
    )
    conn.commit()
    conn.close()
    try: log_audit(None, "document_signed", f"Token: {token[:8]}..., IP: {ip}")
    except: pass
    return {"message": "Document signed successfully"}


# ==================== DOCUMENT VIEWER ANALYTICS ====================

@app.get("/api/documents/{doc_id}/analytics")
async def get_document_analytics(doc_id: int, user=Depends(get_current_user)):
    """Get analytics for a document"""
    conn = get_db()
    doc = conn.execute("SELECT user_id FROM documents WHERE id = ?", (doc_id,)).fetchone()
    if not doc or doc["user_id"] != user["id"]:
        conn.close()
        raise HTTPException(status_code=404, detail="Document not found")

    views = conn.execute(
        "SELECT id, viewer_ip, viewer_email, duration_seconds, scroll_depth, opened_at FROM document_views WHERE document_id = ? ORDER BY opened_at DESC LIMIT 50",
        (doc_id,)
    ).fetchall()
    total = conn.execute("SELECT COUNT(*) as c FROM document_views WHERE document_id = ?", (doc_id,)).fetchone()["c"]
    avg_duration = conn.execute("SELECT AVG(duration_seconds) as avg FROM document_views WHERE document_id = ? AND duration_seconds > 0", (doc_id,)).fetchone()["avg"]
    avg_scroll = conn.execute("SELECT AVG(scroll_depth) as avg FROM document_views WHERE document_id = ? AND scroll_depth > 0", (doc_id,)).fetchone()["avg"]

    conn.close()
    return {
        "total_views": total,
        "avg_duration_seconds": round(avg_duration or 0, 1),
        "avg_scroll_depth": round((avg_scroll or 0) * 100, 1),
        "views": [dict(v) for v in views]
    }

@app.post("/api/track/view")
async def track_document_view(request_data: dict = Body(...), request: Request = None):
    """Track a document view (called from published/shared pages)"""
    doc_id = request_data.get("document_id")
    duration = request_data.get("duration_seconds", 0)
    scroll_depth = request_data.get("scroll_depth", 0)
    viewer_email = request_data.get("viewer_email")

    if not doc_id:
        raise HTTPException(status_code=400, detail="document_id required")

    ip = request.client.host if request and request.client else "unknown"
    conn = get_db()
    conn.execute(
        "INSERT INTO document_views (document_id, viewer_ip, viewer_email, duration_seconds, scroll_depth) VALUES (?, ?, ?, ?, ?)",
        (doc_id, ip, viewer_email, duration, scroll_depth)
    )
    conn.commit()
    conn.close()
    return {"message": "View tracked"}


# ==================== HEALTH CHECK ====================

@app.get("/api/health")
async def health_check():
    """Check system health"""
    status = {"api": "ok", "database": "ok", "ollama": "unknown", "gemini": "unknown"}

    # Check DB
    try:
        conn = get_db()
        conn.execute("SELECT 1")
        conn.close()
    except Exception:
        status["database"] = "error"

    # Check Ollama
    try:
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(_executor, ollama_client.list)
        status["ollama"] = "running"
    except Exception:
        status["ollama"] = "not running"

    # Check Gemini key
    api_key = GEMINI_API_KEY
    if not api_key:
        conn = get_db()
        row = conn.execute("SELECT value FROM settings WHERE key = 'gemini_api_key'").fetchone()
        conn.close()
        if row:
            api_key = row["value"]
    status["gemini"] = "configured" if api_key else "no api key"

    return status


# ==================== PRESENTATION / SLIDE GENERATION ====================

@app.post("/api/generate/presentation")
async def generate_presentation(request_data: dict = Body(...), user=Depends(get_optional_user)):
    """Generate a PowerPoint presentation from a prompt"""
    prompt = request_data.get("prompt", "")
    slide_count = request_data.get("slide_count", 8)
    model = request_data.get("model", "gemini-3-flash-preview")

    if not prompt:
        raise HTTPException(status_code=400, detail="Prompt required")

    # Ask AI to generate slide structure as JSON
    slide_prompt = f"""Create a {slide_count}-slide presentation about: {prompt}

Return ONLY a JSON array of slides. Each slide object must have:
- "title": string (slide title)
- "content": array of strings (bullet points, 3-5 per slide)
- "notes": string (speaker notes, 1-2 sentences)
- "layout": one of "title", "content", "two_column", "section", "quote"

Slide 1 should be layout "title" with the presentation title.
Last slide should be layout "section" as a thank you/Q&A slide.

Return ONLY the JSON array, no markdown formatting."""

    try:
        result = await generate_content(slide_prompt, model_id=model, temperature=0.7, max_tokens=3000)
        clean = result.strip()
        if clean.startswith("```"):
            clean = clean.split("\n", 1)[1].rsplit("```", 1)[0]
        slides = json.loads(clean)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate slides: {str(e)}")

    # Generate PPTX
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides):
        layout_name = slide_data.get("layout", "content")
        title_text = slide_data.get("title", f"Slide {i+1}")
        content_items = slide_data.get("content", [])
        notes_text = slide_data.get("notes", "")

        if layout_name == "title":
            layout = prs.slide_layouts[0]  # Title slide
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title_text
            if slide.placeholders[1]:
                slide.placeholders[1].text = "\n".join(content_items) if content_items else ""
        elif layout_name == "section":
            layout = prs.slide_layouts[2]  # Section header
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title_text
        else:
            layout = prs.slide_layouts[1]  # Title and content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title_text
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for j, item in enumerate(content_items):
                    if j == 0:
                        tf.paragraphs[0].text = item
                    else:
                        p = tf.add_paragraph()
                        p.text = item
                    tf.paragraphs[j].level = 0

        # Add speaker notes
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    pptx_bytes = buffer.read()

    # Save to output dir
    timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    filename = f"Presentation_{timestamp}.pptx"
    save_path = os.path.join(OUTPUT_DIR, filename)
    with open(save_path, "wb") as f:
        f.write(pptx_bytes)

    # Return as downloadable
    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )


@app.post("/api/generate/presentation/preview")
async def preview_presentation(request_data: dict = Body(...), user=Depends(get_optional_user)):
    """Generate presentation structure as JSON (for preview)"""
    prompt = request_data.get("prompt", "")
    slide_count = request_data.get("slide_count", 8)
    model = request_data.get("model", "gemini-3-flash-preview")

    if not prompt:
        raise HTTPException(status_code=400, detail="Prompt required")

    slide_prompt = f"""Create a {slide_count}-slide presentation about: {prompt}

Return ONLY a JSON array of slides. Each slide object must have:
- "title": string
- "content": array of strings (bullet points)
- "notes": string (speaker notes)
- "layout": one of "title", "content", "two_column", "section", "quote"

Return ONLY the JSON array, no markdown."""

    result = await generate_content(slide_prompt, model_id=model, temperature=0.7, max_tokens=3000)
    clean = result.strip()
    if clean.startswith("```"):
        clean = clean.split("\n", 1)[1].rsplit("```", 1)[0]
    slides = json.loads(clean)

    return {"slides": slides, "count": len(slides)}


# ==================== BULK BATCH GENERATION ====================

@app.post("/api/generate/batch")
async def batch_generate(
    file: UploadFile = File(...),
    skill_name: str = Form(None),
    prompt_template: str = Form(...),
    model: str = Form("gemini-3-flash-preview"),
    user=Depends(get_current_user)
):
    """Bulk generate documents from CSV data. Use {{column_name}} as merge fields in prompt_template."""
    import pandas as pd
    import zipfile

    # Read CSV
    content = await file.read()
    try:
        df = pd.read_csv(io.BytesIO(content))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Invalid CSV: {str(e)}")

    if len(df) > 100:
        raise HTTPException(status_code=400, detail="Maximum 100 rows per batch")
    if len(df) == 0:
        raise HTTPException(status_code=400, detail="CSV is empty")

    # Generate documents
    results = []
    errors = []

    for idx, row in df.iterrows():
        try:
            # Replace merge fields in template
            filled_prompt = prompt_template
            for col in df.columns:
                filled_prompt = filled_prompt.replace(f"{{{{{col}}}}}", str(row[col]))

            # Generate
            result = await generate_content(filled_prompt, model_id=model, temperature=0.7, max_tokens=2000)
            results.append({"row": idx + 1, "content": result, "status": "success"})
        except Exception as e:
            errors.append({"row": idx + 1, "error": str(e)})
            results.append({"row": idx + 1, "content": f"Error: {str(e)}", "status": "error"})

    # Create ZIP with all documents
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
        for r in results:
            if r["status"] == "success":
                filename = f"document_row_{r['row']}.md"
                zf.writestr(filename, r["content"])

        # Include summary
        summary = f"Batch Generation Summary\n{'='*40}\n"
        summary += f"Total rows: {len(df)}\n"
        summary += f"Successful: {sum(1 for r in results if r['status']=='success')}\n"
        summary += f"Errors: {len(errors)}\n"
        if errors:
            summary += f"\nErrors:\n"
            for e in errors:
                summary += f"  Row {e['row']}: {e['error']}\n"
        zf.writestr("_summary.txt", summary)

    zip_buffer.seek(0)

    timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    return Response(
        content=zip_buffer.read(),
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="batch_{timestamp}.zip"'}
    )


@app.post("/api/batch/preview")
async def preview_batch_csv(file: UploadFile = File(...)):
    """Preview CSV columns and first 3 rows"""
    import pandas as pd
    content = await file.read()
    try:
        df = pd.read_csv(io.BytesIO(content))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Invalid CSV: {str(e)}")

    return {
        "columns": list(df.columns),
        "row_count": len(df),
        "preview_rows": df.head(3).to_dict(orient="records"),
        "merge_fields": [f"{{{{{col}}}}}" for col in df.columns]
    }


# ==================== STATIC FILE SERVING (Desktop Mode) ====================

DIST_DIR = os.path.join(os.path.dirname(__file__), "..", "dist")

@app.get("/p/{slug}")
async def view_published_document(slug: str, request: Request):
    conn = get_db()
    doc = conn.execute("SELECT * FROM documents WHERE publish_slug = ? AND is_published = 1", (slug,)).fetchone()
    conn.close()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    # Log view
    try: log_audit(None, "doc_view", f"Slug: {slug}, IP: {request.client.host}")
    except: pass
    html_content = doc["html_content"] or markdown_to_html(doc["content"])
    return HTMLResponse(f"""<!DOCTYPE html><html><head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>{doc["title"]}</title>
<style>body{{font-family:'Segoe UI',system-ui,sans-serif;max-width:800px;margin:40px auto;padding:20px;line-height:1.6;color:#1a1a2e}}
h1{{color:#0f172a;border-bottom:2px solid #e2e8f0;padding-bottom:8px}}h2{{color:#1e293b}}h3{{color:#334155}}
table{{border-collapse:collapse;width:100%;margin:16px 0}}th,td{{border:1px solid #e2e8f0;padding:8px 12px}}
th{{background:#f8fafc}}code{{background:#f1f5f9;padding:2px 6px;border-radius:4px}}
blockquote{{border-left:4px solid #e2e8f0;margin:16px 0;padding:8px 16px;color:#64748b}}
.footer{{text-align:center;margin-top:40px;padding-top:20px;border-top:1px solid #e2e8f0;color:#94a3b8;font-size:12px}}</style>
</head><body>{html_content}<div class="footer">Published with Universal Document Creator</div></body></html>""")

if os.path.isdir(DIST_DIR):
    # Serve static assets (JS, CSS, images)
    app.mount("/assets", StaticFiles(directory=os.path.join(DIST_DIR, "assets")), name="static-assets")

    @app.get("/{full_path:path}")
    async def serve_spa(full_path: str):
        """Serve the SPA — any non-API route returns index.html.
        Unknown /api/ routes get a proper 404 JSON response."""
        # Don't serve SPA for unknown API routes — return proper 404
        if full_path.startswith("api/"):
            raise HTTPException(status_code=404, detail="API endpoint not found")
        file_path = os.path.join(DIST_DIR, full_path)
        if full_path and os.path.isfile(file_path):
            return FileResponse(file_path)
        # SPA fallback — return index.html for all routes
        index_path = os.path.join(DIST_DIR, "index.html")
        if os.path.isfile(index_path):
            return FileResponse(index_path)
        raise HTTPException(status_code=404, detail="Not found")

# ==================== MAIN ====================

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)
